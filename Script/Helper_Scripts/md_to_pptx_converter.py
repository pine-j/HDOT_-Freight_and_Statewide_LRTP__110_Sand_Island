"""
Markdown to PowerPoint Converter

Converts structured markdown files into professional PowerPoint presentations
with support for HDOT branding and custom templates.

Usage:
    python md_to_pptx_converter.py input.md
    python md_to_pptx_converter.py input.md --template template.potx
    python md_to_pptx_converter.py input.md -o output.pptx

Requirements:
    pip install python-pptx

Author: Generated for HDOT presentations
"""

import argparse
import copy
import logging
import os
import re
import shutil
import tempfile
import warnings
import zipfile
from dataclasses import dataclass, field
from enum import Enum
from typing import List, Optional, Dict, Any

from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor

# Configure module-level logger
logger = logging.getLogger(__name__)

# Suppress harmless duplicate name warnings when clearing template slides
warnings.filterwarnings('ignore', message='Duplicate name:.*ppt/slides/')


def _clear_metadata(prs: Presentation) -> None:
    """Clear all core document properties (author, title, dates, etc.)
    so the exported file contains no identifying metadata."""
    props = prs.core_properties
    props.author = ""
    props.last_modified_by = ""
    props.title = ""
    props.subject = ""
    props.keywords = ""
    props.category = ""
    props.comments = ""
    props.content_status = ""
    props.revision = 1
    # Set created / modified to the Unix epoch so they appear blank
    epoch = datetime(1970, 1, 1)
    props.created = epoch
    props.modified = epoch
    logger.debug("Document metadata cleared.")


# ============================================
# CONFIGURATION - Modify these settings
# ============================================

CONFIG = {
    # Official HDOT Brand Colors (from State of Hawaii Web Style Guide)
    # These are ONLY used when template_path is None
    # Using Makai (Ocean) theme as primary - best for general HDOT branding
    "colors": {
        # Makai (Ocean) Theme - Primary
        "primary_blue": (38, 72, 110),      # #26486E - Dark Blue (primary brand color)
        "medium_blue": (77, 118, 164),      # #4D76A4 - Medium Blue (secondary)
        "dark_gray": (85, 85, 85),          # #555555 - Body text, neutral accents
        
        # Standard colors
        "white": (255, 255, 255),           # #ffffff
        "light_gray": (218, 225, 235),      # #DAE1EB - Blue-gray for alternating rows
        
        # Mauka (Mountain) Theme - Alternative for inland projects
        "dark_green": (47, 89, 26),         # #2F591A
        "medium_green": (46, 114, 20),      # #2E7214
        "light_green": (128, 192, 98),      # #80C062
        
        # Aina (Land/Earth) Theme - Alternative for sustainability
        "dark_teal": (23, 84, 80),          # #175450
        "teal": (34, 118, 124),             # #22767C
        "light_teal": (115, 201, 207),      # #73C9CF
        
        # Sandy (Soil/Sand) Theme - Alternative for construction
        "dark_brown": (85, 65, 28),         # #55411C
        "medium_brown": (137, 113, 72),     # #897148
        "beige": (191, 183, 137),           # #BFB789
    },
    
    # Color assignments for slide elements (when NOT using template)
    "color_mapping": {
        "title_slide_bg": "primary_blue",      # Dark Blue #26486E
        "title_slide_text": "white",
        "section_slide_bar": "medium_blue",    # Medium Blue #4D76A4
        "section_slide_text": "white",
        "content_title_bar": "primary_blue",   # Dark Blue #26486E
        "content_title_text": "white",
        "body_text": "dark_gray",              # Dark Gray #555555
        "table_header_bg": "primary_blue",     # Dark Blue #26486E
        "table_header_text": "white",
        "table_alt_row": "light_gray",         # Light Gray #dadee5
        "accent": "medium_blue",               # Medium Blue #4D76A4
        "bullet_color": "medium_blue",          # Medium Blue #4D76A4 — Hawaii Makai (Ocean) theme
    },
    
    # Font settings (HDOT standard: Open Sans)
    "fonts": {
        "title": {"name": "Calibri", "size": 44, "bold": True},
        "section": {"name": "Calibri", "size": 36, "bold": True},
        "slide_title": {"name": "Calibri", "size": 28, "bold": True},
        "body": {"name": "Calibri", "size": 18},
        "table_header": {"name": "Calibri", "size": 14, "bold": True},
        "table_body": {"name": "Calibri", "size": 14},
    },
    
    # Slide dimensions (widescreen 16:9)
    "slide_width": 13.333,  # inches
    "slide_height": 7.5,    # inches
}


# ============================================
# DATA STRUCTURES
# ============================================

class SlideType(Enum):
    """Enumeration of supported slide types"""
    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"


# Placeholder type constants for readability
PH_TITLE = PP_PLACEHOLDER.TITLE           # 1
PH_BODY = PP_PLACEHOLDER.BODY             # 2
PH_CENTER_TITLE = PP_PLACEHOLDER.CENTER_TITLE  # 3
PH_SUBTITLE = PP_PLACEHOLDER.SUBTITLE     # 4
PH_OBJECT = PP_PLACEHOLDER.OBJECT         # 7
PH_PICTURE = PP_PLACEHOLDER.PICTURE       # 18

# Media clip placeholder type (no enum constant in python-pptx)
PH_MEDIA_CLIP = 14

# Slide-number placeholder type constant
PH_SLIDE_NUMBER = PP_PLACEHOLDER.SLIDE_NUMBER  # 13


def _ensure_slide_number(slide):
    """Copy the slide-number placeholder from the layout onto the slide.

    When python-pptx creates a slide via ``add_slide()``, it generates
    ``<p:sp>`` elements for *content* placeholders (title, body) but does
    **not** copy "decorator" placeholders (slide number, footer, date) into
    the slide's shape tree.  Those decorators are expected to be inherited
    from the layout; however, PowerPoint only renders inherited decorators
    when the user has checked *Insert → Header & Footer → Slide number*.

    By deep-copying the layout's slide-number ``<p:sp>`` element directly
    into the slide's ``<p:spTree>``, the ``<a:fld type="slidenum">`` field
    is always present and the number renders unconditionally.

    If the layout has no slide-number placeholder (e.g. Title Slide) or if
    the slide already contains one, this function is a harmless no-op.
    """
    from lxml import etree as _etree

    # Quick check: does the slide already have a sldNum placeholder?
    spTree = slide._element.find(qn('p:cSld'))
    if spTree is None:
        return
    spTree = spTree.find(qn('p:spTree'))
    if spTree is None:
        return

    for sp in spTree.iter(qn('p:sp')):
        nvPr = sp.find(qn('p:nvSpPr'))
        if nvPr is not None:
            ph_el = nvPr.find(qn('p:nvPr'))
            if ph_el is not None:
                ph = ph_el.find(qn('p:ph'))
                if ph is not None and ph.get('type') == 'sldNum':
                    return  # Already present

    # Find the slide-number placeholder on the parent layout
    layout = slide.slide_layout
    src_sp = None
    for shape in layout.placeholders:
        if shape.placeholder_format.type == PH_SLIDE_NUMBER:
            src_sp = shape._element
            break

    if src_sp is None:
        return  # Layout has no slide-number placeholder (e.g. Title Slide)

    # Deep-copy the XML element and append it to the slide's shape tree
    sp_copy = copy.deepcopy(src_sp)
    spTree.append(sp_copy)
    logger.debug("    -> Copied slide-number placeholder from layout")


@dataclass
class TextRun:
    """Represents a formatted text run"""
    text: str
    bold: bool = False
    italic: bool = False
    link: Optional[str] = None


@dataclass
class BulletItem:
    """Represents a bullet point with optional nesting"""
    text_runs: List[TextRun]
    level: int = 0  # 0 = top level, 1 = nested, etc.
    is_ordered: bool = False  # True for numbered list items
    is_label: bool = False  # True = render as non-bulleted label/paragraph


@dataclass
class TableData:
    """Represents a table"""
    headers: List[str]
    rows: List[List[str]]


@dataclass
class ContentBlock:
    """A block of content - either bullets or a table.
    
    Used to preserve the original order of interleaved bullets and tables
    within a single slide.
    """
    block_type: str  # "bullets" or "table"
    bullets: List[BulletItem] = field(default_factory=list)
    table: Optional[TableData] = None


@dataclass
class SlideContent:
    """Represents the content of a slide"""
    slide_type: SlideType
    title: str
    subtitle: str = ""
    content_blocks: List[ContentBlock] = field(default_factory=list)
    has_image: bool = False  # For future image support


@dataclass
class LayoutInfo:
    """Information about a slide layout"""
    index: int
    name: str
    has_title_placeholder: bool = False
    has_body_placeholder: bool = False
    has_subtitle_placeholder: bool = False
    has_picture_placeholder: bool = False
    placeholder_count: int = 0


# ============================================
# MARKDOWN PARSER
# ============================================

class MarkdownToSlides:
    """Parse markdown and extract slide structure"""
    
    def __init__(self, markdown_text: str):
        self.markdown_text = markdown_text
        self.slides: List[SlideContent] = []
        self.current_h1 = None
        self.current_h2 = None
        self.current_h3 = None
        self.current_slide_content = []
        self._in_html_block = False
        
    def parse(self) -> List[SlideContent]:
        """Parse markdown and return list of slides"""
        lines = self.markdown_text.split('\n')
        i = 0
        
        while i < len(lines):
            line = lines[i]
            
            # Handle multi-line HTML/CSS blocks (e.g., <style>...</style>)
            if self._in_html_block:
                if '</style>' in line or '</script>' in line:
                    self._in_html_block = False
                i += 1
                continue
            
            # Detect start of multi-line HTML blocks
            stripped_line = line.strip()
            if stripped_line.startswith('<style') or stripped_line.startswith('<script'):
                self._in_html_block = True
                # Check if it closes on the same line
                if '</style>' in line or '</script>' in line:
                    self._in_html_block = False
                i += 1
                continue
            
            # Skip standalone HTML comment lines (<!-- ... -->)
            if stripped_line.startswith('<!--') and stripped_line.endswith('-->'):
                i += 1
                continue
            
            # H1 - Title slide
            if line.startswith('# ') and not line.startswith('## '):
                self._flush_current_slide()
                title = line[2:].strip()
                # Check if next line is subtitle (non-empty, non-heading, non-bullet)
                subtitle = ""
                if i + 1 < len(lines):
                    next_line = lines[i + 1].strip()
                    if (next_line 
                            and not next_line.startswith('#')
                            and not next_line.startswith('- ')
                            and not next_line.startswith('* ')
                            and not re.match(r'^\d+\.\s', next_line)
                            and not next_line.startswith('|')):
                        subtitle = next_line
                        i += 1
                
                self.slides.append(SlideContent(
                    slide_type=SlideType.TITLE,
                    title=title,
                    subtitle=subtitle
                ))
                self.current_h1 = title
                i += 1
                continue
            
            # H2 - Section divider
            if line.startswith('## ') and not line.startswith('### '):
                self._flush_current_slide()
                title = line[3:].strip()
                self.slides.append(SlideContent(
                    slide_type=SlideType.SECTION,
                    title=title
                ))
                self.current_h2 = title
                self.current_h3 = None
                i += 1
                continue
            
            # H3 - Content slide
            if line.startswith('### '):
                self._flush_current_slide()
                self.current_h3 = line[4:].strip()
                i += 1
                continue
            
            # Horizontal rule - slide break
            if stripped_line == '---':
                self._flush_current_slide()
                # Keep current H3 for continuation
                i += 1
                continue
            
            # Table detection
            if '|' in line and stripped_line.startswith('|'):
                table_lines = [line]
                i += 1
                # Collect all table lines (must also start with |)
                while i < len(lines) and '|' in lines[i] and lines[i].strip().startswith('|'):
                    table_lines.append(lines[i])
                    i += 1
                
                table = self._parse_table(table_lines)
                if table:
                    self.current_slide_content.append(('table', table))
                continue
            
            # Bullet points (unordered)
            if stripped_line.startswith('- ') or stripped_line.startswith('* '):
                bullet = self._parse_bullet(line)
                if bullet:
                    self.current_slide_content.append(('bullet', bullet))
                i += 1
                continue
            
            # Numbered lists (ordered)
            numbered_match = re.match(r'^(\s*)(\d+)\.\s+(.*)', line)
            if numbered_match:
                bullet = self._parse_numbered_item(line, numbered_match)
                if bullet:
                    self.current_slide_content.append(('bullet', bullet))
                i += 1
                continue
            
            # Plain text paragraph (non-empty, non-heading, non-bullet, non-table).
            # Rendered as a non-bulleted paragraph.  Useful for section labels
            # like "**Filtering Logic:**" written without a bullet marker.
            if stripped_line and self.current_h3:
                text_runs = self._parse_inline_formatting(stripped_line)
                self.current_slide_content.append(('bullet', BulletItem(
                    text_runs=text_runs, level=0, is_label=True
                )))
                i += 1
                continue
            
            # Empty lines or other content
            i += 1
        
        # Flush any remaining content
        self._flush_current_slide()
        
        return self.slides
    
    def _flush_current_slide(self):
        """Create a slide from accumulated content.
        
        Groups consecutive bullets into a single ContentBlock and wraps each
        table in its own ContentBlock, preserving the original interleaved order.
        """
        if self.current_h3 and self.current_slide_content:
            content_blocks = []
            current_bullets = []
            
            for item_type, item_data in self.current_slide_content:
                if item_type == 'bullet':
                    current_bullets.append(item_data)
                elif item_type == 'table':
                    # Flush any accumulated bullets before adding table
                    if current_bullets:
                        content_blocks.append(ContentBlock(
                            block_type="bullets",
                            bullets=current_bullets
                        ))
                        current_bullets = []
                    content_blocks.append(ContentBlock(
                        block_type="table",
                        table=item_data
                    ))
            
            # Flush any remaining bullets
            if current_bullets:
                content_blocks.append(ContentBlock(
                    block_type="bullets",
                    bullets=current_bullets
                ))
            
            slide = SlideContent(
                slide_type=SlideType.CONTENT,
                title=self.current_h3,
                content_blocks=content_blocks
            )
            self.slides.append(slide)
            self.current_slide_content = []
    
    def _parse_bullet(self, line: str) -> Optional[BulletItem]:
        """Parse a bullet point line.
        
        Auto-detects "label" lines: bullets whose text is ONLY bold
        (optionally ending with a colon), e.g. ``- **Filtering Logic:**``.
        These are rendered without a bullet character as section labels.
        """
        # Determine nesting level by leading spaces
        stripped = line.lstrip()
        if not (stripped.startswith('- ') or stripped.startswith('* ')):
            return None
        
        indent = len(line) - len(stripped)
        level = indent // 2  # 2 spaces per level
        
        # Remove bullet marker
        text = stripped[2:].strip()
        
        # Auto-detect label lines: text is ONLY bold, optionally with colon
        # Examples that match:  **Filtering Logic:**   **Key Assumptions**:
        # Examples that DON'T:  **Note:** Some text    **bold** and normal
        is_label = bool(re.match(r'^\*\*[^*]+\*\*:?\s*$', text))
        
        # Parse inline formatting
        text_runs = self._parse_inline_formatting(text)
        
        return BulletItem(text_runs=text_runs, level=min(level, 2), is_label=is_label)
    
    def _parse_numbered_item(self, line: str, match: re.Match) -> Optional[BulletItem]:
        """Parse a numbered list item (e.g., '1. Item text').
        
        Preserves the original numbering by prepending it as a bold text run
        (e.g., '1. ') before the item text.
        """
        indent_str = match.group(1)
        number = match.group(2)
        text = match.group(3).strip()
        
        indent = len(indent_str)
        level = indent // 2  # 2 spaces per level
        
        # Parse inline formatting for the item text
        text_runs = self._parse_inline_formatting(text)
        
        # Prepend the number as a bold run to preserve ordering
        number_run = TextRun(text=f"{number}. ", bold=True)
        text_runs.insert(0, number_run)
        
        return BulletItem(text_runs=text_runs, level=min(level, 2), is_ordered=True)
    
    def _parse_inline_formatting(self, text: str) -> List[TextRun]:
        """Parse bold, italic, and links in text"""
        runs = []
        current_pos = 0
        
        # Pattern for **bold**, *italic*, and [link](url)
        pattern = r'(\*\*.*?\*\*|\*.*?\*|\[.*?\]\(.*?\))'
        
        for match in re.finditer(pattern, text):
            # Add text before match
            if match.start() > current_pos:
                runs.append(TextRun(text=text[current_pos:match.start()]))
            
            matched_text = match.group(1)
            
            # Bold
            if matched_text.startswith('**') and matched_text.endswith('**'):
                runs.append(TextRun(text=matched_text[2:-2], bold=True))
            # Italic
            elif matched_text.startswith('*') and matched_text.endswith('*'):
                runs.append(TextRun(text=matched_text[1:-1], italic=True))
            # Link
            elif matched_text.startswith('['):
                link_match = re.match(r'\[(.*?)\]\((.*?)\)', matched_text)
                if link_match:
                    link_text, url = link_match.groups()
                    runs.append(TextRun(text=link_text, link=url))
            
            current_pos = match.end()
        
        # Add remaining text
        if current_pos < len(text):
            runs.append(TextRun(text=text[current_pos:]))
        
        return runs if runs else [TextRun(text=text)]
    
    def _parse_table(self, lines: List[str]) -> Optional[TableData]:
        """Parse markdown table"""
        if len(lines) < 2:
            return None
        
        # Parse header
        header_line = lines[0].strip()
        headers = [cell.strip() for cell in header_line.split('|')[1:-1]]
        
        if not headers:
            return None
        
        num_cols = len(headers)
        
        # Skip separator line (line 1)
        # Parse data rows
        rows = []
        for line in lines[2:]:
            if not line.strip():
                continue
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            if cells:
                # Normalize row to match header column count
                if len(cells) < num_cols:
                    cells.extend([""] * (num_cols - len(cells)))
                elif len(cells) > num_cols:
                    cells = cells[:num_cols]
                rows.append(cells)
        
        return TableData(headers=headers, rows=rows)


# ============================================
# LAYOUT MANAGER - Intelligent Layout Selection
# ============================================

class LayoutManager:
    """
    Manages intelligent layout selection from PowerPoint templates.
    Analyzes available layouts and maps content types to appropriate layouts.
    """
    
    # Keywords to identify layout types by name
    LAYOUT_KEYWORDS = {
        "title": ["title", "cover", "opening"],
        "section": ["section", "divider", "header"],
        "content": ["content", "body", "text", "bullet"],
        "image": ["image", "picture", "photo", "media"],
        "table": ["table", "comparison"],
        "blank": ["blank", "empty"],
        "two_column": ["two", "column", "side"],
    }
    
    def __init__(self, presentation: 'Presentation'):
        self.prs = presentation
        self.layouts: List[LayoutInfo] = []
        self.layout_map: Dict[str, int] = {}
        self._analyze_layouts()
        self._build_layout_map()
    
    def _analyze_layouts(self):
        """Analyze all layouts in the presentation template"""
        for idx, layout in enumerate(self.prs.slide_layouts):
            info = LayoutInfo(
                index=idx,
                name=layout.name or f"Layout {idx}",
                placeholder_count=0
            )
            
            # Analyze placeholders in the layout
            placeholder_details = []
            for shape in layout.placeholders:
                info.placeholder_count += 1
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                
                # Collect details for debugging
                placeholder_details.append(f"idx={ph_idx}:type={ph_type}")
                
                # Check placeholder types using enum constants
                if ph_type in [PH_TITLE, PH_CENTER_TITLE]:
                    info.has_title_placeholder = True
                elif ph_type in [PH_BODY, PH_OBJECT]:
                    info.has_body_placeholder = True
                elif ph_type == PH_SUBTITLE:
                    info.has_subtitle_placeholder = True
                elif ph_type in [PH_PICTURE, PH_MEDIA_CLIP]:
                    info.has_picture_placeholder = True
                
                # Also check by index for common patterns
                # Index 0 is usually title, Index 1 is usually body/content
                if ph_idx == 1 and not info.has_body_placeholder:
                    # Index 1 often holds content even if type is unusual
                    info.has_body_placeholder = True
            
            self.layouts.append(info)
            logger.debug(
                "  Layout %d: '%s' - Title:%s, Body:%s, Subtitle:%s, Picture:%s [%s]",
                idx, info.name, info.has_title_placeholder,
                info.has_body_placeholder, info.has_subtitle_placeholder,
                info.has_picture_placeholder, ', '.join(placeholder_details)
            )
    
    def _build_layout_map(self):
        """Build a mapping from content types to best layout indices"""
        
        # Find layouts by analyzing names and placeholders
        for info in self.layouts:
            name_lower = info.name.lower()
            
            # Title Slide - exact match for Sand Island template
            if name_lower == "title slide":
                self.layout_map["title_no_image"] = info.index
                self.layout_map["title_with_image"] = info.index  # Same layout for both
            
            # Transition Slide - for section dividers
            if "transition" in name_lower:
                self.layout_map["section"] = info.index
            
            # Content Slide - for regular content
            if name_lower == "content slide":
                self.layout_map["content"] = info.index
            
            # Agenda or Table of Contents - can be used as alternative section or TOC
            if "agenda" in name_lower or "table of contents" in name_lower:
                if "agenda" not in self.layout_map:
                    self.layout_map["agenda"] = info.index
            
            # Generic fallback matching
            # Title/Cover layouts (with subtitle, without image)
            if any(kw in name_lower for kw in ["cover without", "title slide"]):
                if "title_no_image" not in self.layout_map:
                    self.layout_map["title_no_image"] = info.index
            
            # Title/Cover layouts (with image)
            if any(kw in name_lower for kw in ["cover with image", "title with picture"]):
                if info.has_picture_placeholder:
                    if "title_with_image" not in self.layout_map:
                        self.layout_map["title_with_image"] = info.index
            
            # Section header layouts
            if any(kw in name_lower for kw in ["section", "divider"]):
                if "section" not in self.layout_map:
                    self.layout_map["section"] = info.index
            
            # Content layouts (title + content) - relaxed matching
            if "title and content" in name_lower or ("content" in name_lower and "cover" not in name_lower and "table of contents" not in name_lower):
                if info.has_title_placeholder:
                    if "content" not in self.layout_map:
                        self.layout_map["content"] = info.index
            
            # Blank layout
            if "blank" in name_lower or info.placeholder_count == 0:
                if "blank" not in self.layout_map:
                    self.layout_map["blank"] = info.index
        
        # Fallback mappings based on placeholder analysis if names didn't match
        if "title_no_image" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and info.has_subtitle_placeholder and not info.has_picture_placeholder:
                    self.layout_map["title_no_image"] = info.index
                    break
        
        if "content" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and info.has_body_placeholder:
                    self.layout_map["content"] = info.index
                    break
        
        # If still no content layout, use any layout with a title placeholder
        if "content" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and not info.has_picture_placeholder:
                    # Avoid cover slides
                    if "cover" not in info.name.lower():
                        self.layout_map["content"] = info.index
                        break
        
        # For section slides, use transition or agenda layout if no dedicated section layout
        if "section" not in self.layout_map:
            if "agenda" in self.layout_map:
                self.layout_map["section"] = self.layout_map["agenda"]
            elif "title_no_image" in self.layout_map:
                self.layout_map["section"] = self.layout_map["title_no_image"]
        
        # Last resort: use indices commonly found in templates
        if "blank" not in self.layout_map:
            try:
                if len(self.prs.slide_layouts) > 6:
                    self.layout_map["blank"] = 6
                else:
                    self.layout_map["blank"] = len(self.prs.slide_layouts) - 1
            except Exception:
                self.layout_map["blank"] = 0
        
        logger.debug("Layout mapping: %s", self.layout_map)
    
    def get_layout_for_slide(self, slide_content: SlideContent) -> Any:
        """
        Intelligently select the best layout for the given slide content.
        
        Args:
            slide_content: The content to be placed on the slide
            
        Returns:
            The appropriate slide layout object
        """
        layout_key = None
        
        if slide_content.slide_type == SlideType.TITLE:
            # Title slides - check for image
            if slide_content.has_image:
                layout_key = "title_with_image"
            else:
                layout_key = "title_no_image"
        
        elif slide_content.slide_type == SlideType.SECTION:
            layout_key = "section"
        
        elif slide_content.slide_type == SlideType.CONTENT:
            # Content slides - could be enhanced for tables, images, etc.
            layout_key = "content"
        
        # Get the layout index, with fallback to blank
        if layout_key and layout_key in self.layout_map:
            idx = self.layout_map[layout_key]
            logger.debug("  -> Using layout '%s' for %s slide", self.layouts[idx].name, slide_content.slide_type.value)
            return self.prs.slide_layouts[idx]
        
        # Fallback to blank layout
        idx = self.layout_map.get("blank", 0)
        logger.debug("  -> Fallback to blank layout for %s slide", slide_content.slide_type.value)
        return self.prs.slide_layouts[idx]
    
    def get_blank_layout(self) -> Any:
        """Get the blank layout for manual slide creation"""
        idx = self.layout_map.get("blank", len(self.prs.slide_layouts) - 1)
        return self.prs.slide_layouts[idx]
    
    def has_usable_layouts(self) -> bool:
        """Check if the template has usable layouts beyond blank"""
        return len(self.layout_map) > 1


# ============================================
# POWERPOINT GENERATOR
# ============================================

def get_color(color_name: str) -> RGBColor:
    """Get RGBColor from color name in config"""
    rgb = CONFIG["colors"][color_name]
    return RGBColor(*rgb)


def _get_fallback_blank_layout(prs: Presentation):
    """Get a blank layout from a presentation without a LayoutManager.
    
    Falls back to index 6 (common blank index) or the last available layout.
    """
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[len(prs.slide_layouts) - 1]


def convert_potx_to_pptx(potx_path: str) -> str:
    """Convert .potx template to .pptx by modifying content types.
    
    Returns:
        Path to a temporary .pptx file. Caller is responsible for cleanup.
    """
    temp_pptx = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    temp_pptx.close()
    
    # Copy the potx file
    shutil.copy2(potx_path, temp_pptx.name)
    
    # Modify the [Content_Types].xml to change template to presentation
    with zipfile.ZipFile(temp_pptx.name, 'r') as zip_read:
        # Read all files
        file_list = zip_read.namelist()
        
        # Create a new zip with modified content
        temp_zip = tempfile.NamedTemporaryFile(suffix='.zip', delete=False)
        temp_zip.close()
        
        with zipfile.ZipFile(temp_zip.name, 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for file_name in file_list:
                file_data = zip_read.read(file_name)
                
                # Modify [Content_Types].xml
                if file_name == '[Content_Types].xml':
                    file_data = file_data.replace(
                        b'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
                        b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                    )
                
                zip_write.writestr(file_name, file_data)
    
    # Replace the original temp file
    shutil.move(temp_zip.name, temp_pptx.name)
    
    return temp_pptx.name


def create_clean_template(source_template: str, output_path: str) -> str:
    """
    Create a clean template file (no slides) from an existing template.
    
    This is useful when you have a .pptx with slides and want to create
    a clean template that only contains the layouts/masters.
    
    Args:
        source_template: Path to the source template (can have slides)
        output_path: Path for the clean template output
        
    Returns:
        Path to the created clean template
    """
    logger.info("Creating clean template from: %s", source_template)
    
    # Open the source template
    prs = Presentation(source_template)
    
    # Delete all slides
    num_slides = len(prs.slides)
    if num_slides > 0:
        logger.info("Removing %d slides from template...", num_slides)
        for i in range(num_slides - 1, -1, -1):
            try:
                delete_slide(prs, i)
            except Exception:
                pass
    
    # Clear metadata and save as the clean template
    _clear_metadata(prs)
    prs.save(output_path)
    logger.info("Clean template saved to: %s", output_path)
    logger.info("TIP: You can now use this clean template to avoid duplicate warnings.")
    
    return output_path


def delete_slide(prs: Presentation, slide_index: int):
    """
    Delete a slide from the presentation by index.
    Properly removes both the slide reference and the slide part.
    
    Args:
        prs: Presentation object
        slide_index: Index of slide to delete (0-based)
    """
    # Get the slide's relationship ID
    slide_id = prs.slides._sldIdLst[slide_index].rId
    
    # Get the slide part before dropping the relationship
    slide_part = prs.part.related_part(slide_id)
    
    # Remove the slide ID from the presentation's slide list
    del prs.slides._sldIdLst[slide_index]
    
    # Drop the relationship (this disconnects the slide)
    prs.part.drop_rel(slide_id)
    
    # Try to remove the slide part from the package
    try:
        # Access the package and remove the part
        package = prs.part.package
        if hasattr(package, '_parts'):
            # Remove from parts collection if it exists
            parts_to_remove = [p for p in package._parts if p is slide_part]
            for p in parts_to_remove:
                package._parts.remove(p)
    except Exception:
        pass  # If we can't remove the part, it's okay - file will still work


def clear_template_slides(prs: Presentation):
    """
    Remove all existing slides from a template presentation.
    
    This properly deletes slides to avoid duplicate file warnings.
    """
    num_slides = len(prs.slides)
    if num_slides == 0:
        return
    
    logger.info("Clearing %d existing template slides...", num_slides)
    
    # Delete slides from the end to avoid index shifting issues
    deleted_count = 0
    for i in range(num_slides - 1, -1, -1):
        try:
            delete_slide(prs, i)
            deleted_count += 1
        except Exception as e:
            logger.warning("    Could not delete slide %d: %s", i + 1, e)
    
    logger.info("Template slides cleared (%d slides removed).", deleted_count)


def create_presentation(template_path: Optional[str] = None) -> tuple:
    """
    Initialize presentation with template or default settings.
    
    Returns:
        tuple: (Presentation, LayoutManager or None)
    """
    temp_file = None
    layout_manager = None
    
    if template_path and os.path.exists(template_path):
        # Check if it's a .potx file
        if template_path.lower().endswith('.potx'):
            logger.info("Note: .potx template detected. Converting to .pptx format...")
            temp_file = convert_potx_to_pptx(template_path)
            
            try:
                prs = Presentation(temp_file)
                logger.info("Using template: %s", template_path)
            except Exception:
                # Clean up temp file on failure, then re-raise
                if temp_file:
                    try:
                        os.unlink(temp_file)
                    except OSError:
                        pass
                raise
            finally:
                # Clean up temp file after presentation is loaded into memory
                if temp_file:
                    try:
                        os.unlink(temp_file)
                    except OSError:
                        pass
        else:
            prs = Presentation(template_path)
            logger.info("Using template: %s", template_path)
        
        # Clear any existing slides from the template
        if len(prs.slides) > 0:
            clear_template_slides(prs)
        
        # Initialize layout manager for template analysis
        logger.info("Analyzing template layouts...")
        layout_manager = LayoutManager(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(CONFIG["slide_width"])
        prs.slide_height = Inches(CONFIG["slide_height"])
        logger.info("Using default HDOT styling (no template)")
    
    return prs, layout_manager


def _safe_get_first_run(paragraph):
    """Safely get or create the first run of a paragraph after setting text.
    
    After setting `p.text = ...`, python-pptx may or may not create a run object.
    This helper ensures we always have a valid run to format.
    
    Returns:
        The first run of the paragraph.
    """
    if paragraph.runs:
        return paragraph.runs[0]
    # No run was created; add one manually
    run = paragraph.add_run()
    run.text = paragraph.text
    # Clear the direct text so it doesn't duplicate
    paragraph.text = ""
    # Re-fetch since add_run changed the structure
    return paragraph.runs[0] if paragraph.runs else run


def _extract_layout_formatting(shape, slide_layout):
    """Extract font formatting from the matching layout placeholder.
    
    Searches the layout placeholder's runs and lstStyle/defRPr for
    font color, size, and bold settings.  Returns a dict with keys
    'color', 'size', and 'bold' (values may be None if not found).
    """
    result = {"color": None, "size": None, "bold": None}
    ph_idx = shape.placeholder_format.idx
    ph_type = shape.placeholder_format.type
    
    for layout_ph in slide_layout.placeholders:
        # Match by idx, or by type for special placeholders like ctrTitle
        if (layout_ph.placeholder_format.idx == ph_idx
                or layout_ph.placeholder_format.type == ph_type):
            layout_tf = layout_ph.text_frame
            
            # 1. Try run-level formatting (e.g. Title Slide CENTER_TITLE)
            if layout_tf.paragraphs and layout_tf.paragraphs[0].runs:
                run = layout_tf.paragraphs[0].runs[0]
                try:
                    if run.font.color and run.font.color.rgb:
                        result["color"] = run.font.color.rgb
                except Exception:
                    pass
                if run.font.size:
                    result["size"] = run.font.size
                if run.font.bold is not None:
                    result["bold"] = run.font.bold
            
            # 2. Also check lstStyle > defRPr (e.g. Content Slide TITLE)
            lstStyle = layout_ph._element.find(qn('a:lstStyle'))
            if lstStyle is not None:
                defRPr = lstStyle.find('.//' + qn('a:defRPr'))
                if defRPr is not None:
                    if result["size"] is None and defRPr.get('sz'):
                        result["size"] = Pt(int(defRPr.get('sz')) / 100)
                    if result["bold"] is None and defRPr.get('b'):
                        result["bold"] = defRPr.get('b') == '1'
                    if result["color"] is None:
                        fill = defRPr.find(qn('a:solidFill'))
                        if fill is not None:
                            srgb = fill.find(qn('a:srgbClr'))
                            if srgb is not None and srgb.get('val'):
                                result["color"] = RGBColor.from_string(srgb.get('val'))
            break
    
    return result


def _set_placeholder_text(shape, text, slide_layout=None, font_cfg_key="slide_title"):
    """Set text in a template placeholder with proper formatting.
    
    Extracts the color (and optionally size/bold) from the matching layout
    placeholder, then sets the text with the CONFIG font (e.g. Calibri)
    and the template's color.  Also copies explicit position/size from the
    layout so the slide shape doesn't rely on inheritance (which can break
    for special placeholder indices like CENTER_TITLE).
    Enables auto-fit so long titles shrink.
    
    Args:
        shape: The slide placeholder shape to populate.
        text: The text string to set.
        slide_layout: The SlideLayout object for formatting extraction.
        font_cfg_key: Key into CONFIG["fonts"] for font name/size defaults
                      (e.g. "title", "section", "slide_title").
    """
    from lxml import etree
    
    # Extract formatting from the layout placeholder
    fmt = {"color": None, "size": None, "bold": None}
    if slide_layout is not None:
        fmt = _extract_layout_formatting(shape, slide_layout)
        
        # Copy explicit position/size from the layout placeholder to the
        # slide shape.  Without this, the slide shape has empty <p:spPr/>
        # and relies on inheritance, which can fail for special placeholder
        # indices (e.g. CENTER_TITLE idx=4294967295).  Making the geometry
        # explicit ensures the shape has a proper bounding box, appears in
        # thumbnails, and renders at the correct position.
        ph_idx = shape.placeholder_format.idx
        ph_type = shape.placeholder_format.type
        for layout_ph in slide_layout.placeholders:
            if (layout_ph.placeholder_format.idx == ph_idx
                    or layout_ph.placeholder_format.type == ph_type):
                if layout_ph.left is not None:
                    shape.left = layout_ph.left
                    shape.top = layout_ph.top
                    shape.width = layout_ph.width
                    shape.height = layout_ph.height
                    logger.debug("    -> Copied explicit position from layout: "
                                 "left=%.2f\", top=%.2f\", w=%.2f\", h=%.2f\"",
                                 layout_ph.left.inches, layout_ph.top.inches,
                                 layout_ph.width.inches, layout_ph.height.inches)
                break
    
    # Set the text (creates a bare run)
    shape.text = text
    
    # Get the run to format (create one if needed)
    p = shape.text_frame.paragraphs[0]
    run = _safe_get_first_run(p)
    
    # Apply font from CONFIG
    font_cfg = CONFIG["fonts"].get(font_cfg_key, CONFIG["fonts"]["slide_title"])
    run.font.name = font_cfg["name"]
    
    # Apply color from template layout (critical for visual match)
    if fmt["color"]:
        run.font.color.rgb = fmt["color"]
    
    # Apply size: use layout value if found, otherwise CONFIG value
    if fmt["size"]:
        run.font.size = fmt["size"]
    elif "size" in font_cfg:
        run.font.size = Pt(font_cfg["size"])
    
    # Apply bold from layout or CONFIG
    if fmt["bold"] is not None:
        run.font.bold = fmt["bold"]
    elif font_cfg.get("bold"):
        run.font.bold = font_cfg["bold"]
    
    # Enable auto-fit so text shrinks to fit the placeholder
    bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        # Remove any existing autofit settings
        for child in list(bodyPr):
            if child.tag in [qn('a:noAutofit'), qn('a:normAutofit'), qn('a:spAutoFit')]:
                bodyPr.remove(child)
        # Add normAutofit (shrink text on overflow)
        etree.SubElement(bodyPr, qn('a:normAutofit'))


def add_title_slide(prs: Presentation, title: str, subtitle: str = "", 
                    layout_manager: Optional[LayoutManager] = None):
    """
    Add a title slide using template layout or manual creation.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text (optional)
        layout_manager: LayoutManager for intelligent layout selection
    """
    # Try to use template layout
    if layout_manager:
        slide_content = SlideContent(slide_type=SlideType.TITLE, title=title, subtitle=subtitle)
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to populate placeholders
        title_set = False
        subtitle_set = False
        
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            
            # TITLE or CENTER_TITLE
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                _set_placeholder_text(shape, title, slide_layout=slide_layout, font_cfg_key="title")
                title_set = True
                logger.debug("    -> Title set in placeholder type %s", ph_type)
            
            # SUBTITLE
            elif ph_type == PH_SUBTITLE and subtitle and not subtitle_set:
                _set_placeholder_text(shape, subtitle, slide_layout=slide_layout, font_cfg_key="body")
                subtitle_set = True
                logger.debug("    -> Subtitle set in placeholder type %s", ph_type)
        
        if title_set:
            return slide
        
        # If no placeholders found, fall through to manual creation
        logger.debug("    No title placeholders found, using manual layout")
    
    # Fallback: Manual slide creation (original behavior)
    slide_layout = _get_fallback_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, 
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["title_slide_bg"])
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), 
        Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(CONFIG["fonts"]["title"]["size"])
    p.font.bold = CONFIG["fonts"]["title"]["bold"]
    p.font.name = CONFIG["fonts"]["title"]["name"]
    p.font.color.rgb = get_color(CONFIG["color_mapping"]["title_slide_text"])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), 
            Inches(12.333), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.name = CONFIG["fonts"]["title"]["name"]
        p.font.color.rgb = get_color(CONFIG["color_mapping"]["title_slide_text"])
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs: Presentation, title: str,
                      layout_manager: Optional[LayoutManager] = None):
    """
    Add a section divider slide using template layout or manual creation.
    
    Args:
        prs: Presentation object
        title: Section title text
        layout_manager: LayoutManager for intelligent layout selection
    """
    # Try to use template layout
    if layout_manager:
        slide_content = SlideContent(slide_type=SlideType.SECTION, title=title)
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to populate placeholders
        title_set = False
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            # TITLE or CENTER_TITLE
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                _set_placeholder_text(shape, title, slide_layout=slide_layout, font_cfg_key="section")
                title_set = True
                logger.debug("    -> Section title set in placeholder type %s", ph_type)
                break
        
        if title_set:
            return slide
        
        # If no placeholders, try adding text to the title area if layout has one
        logger.debug("    No section placeholders found, using manual layout")
    
    # Fallback: Manual slide creation (original behavior)
    slide_layout = _get_fallback_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3), 
        prs.slide_width, Inches(1.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["section_slide_bar"])
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), 
        Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(CONFIG["fonts"]["section"]["size"])
    p.font.bold = CONFIG["fonts"]["section"]["bold"]
    p.font.name = CONFIG["fonts"]["section"]["name"]
    p.font.color.rgb = get_color(CONFIG["color_mapping"]["section_slide_text"])
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def _populate_placeholder_bullets(shape, bullets: List[BulletItem],
                                   body_format: Optional[Dict[str, Any]] = None):
    """Populate a body placeholder with bullet content.
    
    Unlike add_bullets_to_slide (which creates a manual textbox), this writes
    directly into the template's body placeholder.
    
    Explicit bullet XML (buFont + buChar + buClr) is added because OBJECT-type
    placeholders (common in custom templates) don't inherit bullet styles
    from the slide master the way BODY placeholders do.
    
    When body_format is provided, text color is set explicitly to match
    the textbox-rendered bullets on mixed-content slides, ensuring visual
    consistency across all slides.
    
    Args:
        shape: The body placeholder shape on the slide.
        bullets: List of BulletItem objects to render.
        body_format: Optional dict with 'font_name', 'font_size', 'font_color'
                     extracted from the template body placeholder.
    """
    tf = shape.text_frame
    tf.word_wrap = True
    
    # Clear existing sample / prompt text from the placeholder.
    # Remove all <a:p> elements except the first, then clear runs from the
    # first paragraph while preserving its <a:pPr> (paragraph properties).
    p_elements = tf._txBody.findall(qn('a:p'))
    for p_elem in p_elements[1:]:
        tf._txBody.remove(p_elem)
    first_p = p_elements[0]
    for child in list(first_p):
        if child.tag != qn('a:pPr'):
            first_p.remove(child)
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Set nesting level for potential inheritance from slide master
        p.level = bullet.level
        
        # Add text runs with explicit formatting to match textbox-rendered slides
        for run_data in bullet.text_runs:
            r = p.add_run()
            r.text = run_data.text
            
            # Explicitly set text color for consistency with add_bullets_to_slide
            r.font.color.rgb = (body_format["font_color"]
                                if body_format and body_format.get("font_color")
                                else get_color(CONFIG["color_mapping"]["body_text"]))
            if body_format and body_format.get("font_name"):
                r.font.name = body_format["font_name"]
            if body_format and body_format.get("font_size"):
                r.font.size = body_format["font_size"]
            
            if run_data.bold:
                r.font.bold = True
            if run_data.italic:
                r.font.italic = True
            if run_data.link:
                r.hyperlink.address = run_data.link
        
        # Labels get extra vertical spacing above for visual separation
        if bullet.is_label and i > 0:
            p.space_before = Pt(12)
        
        p.space_after = Pt(8)
        
        # Add explicit bullet formatting — OBJECT placeholders don't inherit
        # bullet styles from the slide master, so we must add buChar XML.
        _add_bullet_formatting(p, bullet.level, bullet.is_ordered, bullet.is_label)
    
    logger.debug("    -> Populated body placeholder with %d bullet items", len(bullets))


def _resolve_scheme_color(prs: Presentation, scheme_name: str) -> Optional[RGBColor]:
    """Resolve a PowerPoint scheme/theme color name to an RGB value.
    
    Looks up the color in the presentation's theme clrScheme.  The theme XML
    is stored in a separate part linked via the slide master's relationships,
    NOT embedded in the master element itself.
    
    Common scheme names: dk1, dk2, lt1, lt2, accent1-6, hlink, folHlink,
    tx1, tx2, bg1, bg2.
    
    Args:
        prs: Presentation object (needed to access the theme)
        scheme_name: Scheme color name (e.g. 'dk1', 'tx1', 'accent1')
    
    Returns:
        RGBColor or None if not resolvable.
    """
    from lxml import etree as _etree
    
    # Map tx/bg aliases to their dk/lt equivalents
    _aliases = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}
    scheme_name = _aliases.get(scheme_name, scheme_name)
    
    try:
        for master in prs.slide_masters:
            # Access the theme through the master's OPC relationships
            for rel in master.part.rels.values():
                if 'theme' in rel.reltype.lower():
                    theme_elem = _etree.fromstring(rel.target_part.blob)
                    clrScheme = theme_elem.find('.//' + qn('a:clrScheme'))
                    if clrScheme is not None:
                        color_elem = clrScheme.find(qn(f'a:{scheme_name}'))
                        if color_elem is not None:
                            srgb = color_elem.find(qn('a:srgbClr'))
                            if srgb is not None and srgb.get('val'):
                                return RGBColor.from_string(srgb.get('val'))
                            sys_clr = color_elem.find(qn('a:sysClr'))
                            if sys_clr is not None and sys_clr.get('lastClr'):
                                return RGBColor.from_string(sys_clr.get('lastClr'))
    except Exception:
        pass
    return None


def _extract_color_from_solidFill(solidFill, prs: Optional[Presentation] = None) -> Optional[RGBColor]:
    """Extract an RGB color from a solidFill element.
    
    Handles both explicit RGB (a:srgbClr) and scheme/theme colors (a:schemeClr).
    """
    if solidFill is None:
        return None
    
    # 1. Direct RGB color
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is not None and srgbClr.get('val'):
        return RGBColor.from_string(srgbClr.get('val'))
    
    # 2. Scheme/theme color — resolve against the presentation theme
    schemeClr = solidFill.find(qn('a:schemeClr'))
    if schemeClr is not None and schemeClr.get('val') and prs is not None:
        resolved = _resolve_scheme_color(prs, schemeClr.get('val'))
        if resolved is not None:
            logger.debug("    Resolved scheme color '%s' -> %s", schemeClr.get('val'), resolved)
            return resolved
    
    return None


def _extract_body_placeholder_format(body_placeholder, slide_layout,
                                     prs: Optional[Presentation] = None) -> Dict[str, Any]:
    """Extract font formatting from a body placeholder before removing it.
    
    Searches the matching layout placeholder's lstStyle (level definitions)
    and run-level formatting for font name, size, and color.
    
    Handles both explicit RGB colors (a:srgbClr) and scheme/theme colors
    (a:schemeClr) by resolving them against the presentation's theme.
    
    Returns:
        Dict with keys 'font_name' (str|None), 'font_size' (Pt|None),
        'font_color' (RGBColor|None).
    """
    fmt: Dict[str, Any] = {"font_name": None, "font_size": None, "font_color": None}
    ph_idx = body_placeholder.placeholder_format.idx
    ph_type = body_placeholder.placeholder_format.type
    
    # Helper: search an element tree for defRPr formatting
    def _extract_from_defRPr(defRPr):
        if defRPr is None:
            return
        if fmt["font_size"] is None and defRPr.get('sz'):
            fmt["font_size"] = Pt(int(defRPr.get('sz')) / 100)
        if fmt["font_name"] is None:
            latin = defRPr.find(qn('a:latin'))
            if latin is not None and latin.get('typeface'):
                fmt["font_name"] = latin.get('typeface')
        if fmt["font_color"] is None:
            solidFill = defRPr.find(qn('a:solidFill'))
            color = _extract_color_from_solidFill(solidFill, prs)
            if color is not None:
                fmt["font_color"] = color
    
    for layout_ph in slide_layout.placeholders:
        if (layout_ph.placeholder_format.idx == ph_idx
                or layout_ph.placeholder_format.type == ph_type):
            # 1. Check lstStyle > lvl1pPr > defRPr (most reliable)
            lstStyle = layout_ph._element.find('.//' + qn('a:lstStyle'))
            if lstStyle is not None:
                lvl1pPr = lstStyle.find(qn('a:lvl1pPr'))
                if lvl1pPr is not None:
                    _extract_from_defRPr(lvl1pPr.find(qn('a:defRPr')))
            
            # 2. Fallback: check run-level formatting in the layout placeholder
            layout_tf = layout_ph.text_frame
            if layout_tf.paragraphs:
                for para in layout_tf.paragraphs:
                    if para.runs:
                        run = para.runs[0]
                        try:
                            if fmt["font_name"] is None and run.font.name:
                                fmt["font_name"] = run.font.name
                        except Exception:
                            pass
                        if fmt["font_size"] is None and run.font.size:
                            fmt["font_size"] = run.font.size
                        try:
                            if fmt["font_color"] is None and run.font.color and run.font.color.rgb:
                                fmt["font_color"] = run.font.color.rgb
                        except Exception:
                            pass
                        break
            break
    
    # 3. If color still missing, check the slide master's body text style
    if fmt["font_color"] is None and prs is not None:
        try:
            master = slide_layout.slide_master
            # Check txStyles > bodyStyle > lvl1pPr > defRPr
            txStyles = master.element.find('.//' + qn('a:txStyles'))
            if txStyles is not None:
                bodyStyle = txStyles.find(qn('a:bodyStyle'))
                if bodyStyle is not None:
                    lvl1pPr = bodyStyle.find(qn('a:lvl1pPr'))
                    if lvl1pPr is not None:
                        defRPr = lvl1pPr.find(qn('a:defRPr'))
                        if defRPr is not None:
                            solidFill = defRPr.find(qn('a:solidFill'))
                            color = _extract_color_from_solidFill(solidFill, prs)
                            if color is not None:
                                fmt["font_color"] = color
                                logger.debug("    -> Got body color from slide master: %s", color)
        except Exception:
            pass
    
    # 4. Final fallback: use the theme's default text color (dk1 / tx1).
    #    This is what PowerPoint uses when no explicit color is specified.
    if fmt["font_color"] is None and prs is not None:
        dk1_color = _resolve_scheme_color(prs, "dk1")
        if dk1_color is not None:
            fmt["font_color"] = dk1_color
            logger.debug("    -> Using theme dk1 color as body text fallback: %s", dk1_color)
    
    logger.debug("    -> Extracted body format: name=%s, size=%s, color=%s",
                 fmt["font_name"], fmt["font_size"], fmt["font_color"])
    return fmt


def _estimate_bullets_height(bullets: List[BulletItem]) -> float:
    """Estimate the vertical space (inches) a bullet block will occupy.

    This is used for *space reservation* (e.g. computing bottom_reserve for
    tables), NOT for rendering.  A slightly tight estimate is preferable to
    an overly generous one — over-estimation starves tables of vertical room
    and produces cramped rows / small fonts.

    Uses PowerPoint's actual ~1.2× line-spacing (not the generous 1.5× used
    in earlier versions) and 6pt inter-item spacing.
    """
    body_size_pt = CONFIG["fonts"]["body"]["size"]
    line_height = body_size_pt * 1.2 / 72   # ~1.2× line-spacing (matches PPT)
    space_after = 6 / 72                     # 6pt inter-item gap, in inches
    text_box_width_inches = 12.333
    avg_char_width_inches = body_size_pt * 0.006  # empirical approximation

    estimated = 0.0
    for bullet in bullets:
        total_chars = sum(len(run.text) for run in bullet.text_runs)
        indent_reduction = bullet.level * 0.5
        effective_chars = max(20, (text_box_width_inches - indent_reduction) / avg_char_width_inches)
        num_lines = max(1, -(-total_chars // int(effective_chars)))  # ceiling division
        estimated += (num_lines * line_height) + space_after

    return max(0.3, estimated)


def add_content_slide(prs: Presentation, slide_content: SlideContent,
                      layout_manager: Optional[LayoutManager] = None):
    """
    Add a content slide using template layout or manual creation.
    
    Renders content blocks (bullets and tables) sequentially, preserving
    their original order from the markdown source.
    
    Args:
        prs: Presentation object
        slide_content: SlideContent with title and content_blocks
        layout_manager: LayoutManager for intelligent layout selection
    """
    slide = None
    current_top = 1.4  # default for manual mode (inches)
    is_template = False
    
    # Try to use template layout
    if layout_manager:
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        title_set = False
        body_placeholder = None
        
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            
            # Title placeholder (TITLE or CENTER_TITLE)
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                _set_placeholder_text(shape, slide_content.title, slide_layout=slide_layout, font_cfg_key="slide_title")
                title_set = True
                logger.debug("    -> Title set in placeholder type %s", ph_type)
            
            # Body/Content placeholder (BODY or OBJECT)
            elif ph_type in [PH_BODY, PH_OBJECT] and body_placeholder is None:
                body_placeholder = shape
        
        if title_set:
            # Determine content composition
            has_tables = any(
                b.block_type == "table" for b in slide_content.content_blocks
            )
            all_bullets = []
            for block in slide_content.content_blocks:
                if block.block_type == "bullets":
                    all_bullets.extend(block.bullets)
            
            # Extract template body formatting (needed by both paths)
            body_format = None
            if body_placeholder:
                body_format = _extract_body_placeholder_format(
                    body_placeholder, slide_layout, prs
                )
            
            if not has_tables and all_bullets and body_placeholder:
                # ── Bullets-only: populate body placeholder directly ──
                # Explicit formatting (text color, bullet color) is applied
                # to ensure consistency with textbox-rendered mixed-content slides.
                _populate_placeholder_bullets(body_placeholder, all_bullets,
                                              body_format=body_format)
                return slide
            
            # ── Mixed content or tables: manual layout ──
            # Remove body placeholder so "Click to add text" doesn't show through
            if body_placeholder:
                try:
                    sp = body_placeholder._element
                    sp.getparent().remove(sp)
                    logger.debug("    -> Removed body placeholder for manual content layout")
                except Exception:
                    pass  # If deletion fails, continue anyway
            is_template = True
            current_top = 1.5  # below template title area
        else:
            logger.debug("    No content placeholders found, using manual layout")
            slide = None  # fall through to manual creation
            body_format = None
    
    if slide is None:
        # Fallback: Manual slide creation (original behavior)
        slide_layout = _get_fallback_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)
        body_format = None
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            prs.slide_width, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["content_title_bar"])
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), 
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_content.title
        p.font.size = Pt(CONFIG["fonts"]["slide_title"]["size"])
        p.font.bold = CONFIG["fonts"]["slide_title"]["bold"]
        p.font.name = CONFIG["fonts"]["slide_title"]["name"]
        p.font.color.rgb = get_color(CONFIG["color_mapping"]["content_title_text"])
        
        current_top = 1.4
    
    # Pre-estimate heights for bullet blocks so tables can reserve space
    # for content that follows them (prevents overflow into footer area).
    block_estimated_heights = []
    for block in slide_content.content_blocks:
        if block.block_type == "bullets" and block.bullets:
            block_estimated_heights.append(_estimate_bullets_height(block.bullets))
        else:
            block_estimated_heights.append(0.0)
    
    # Render content blocks sequentially with tracked vertical position
    gap = 0.15  # inches between blocks
    for idx, block in enumerate(slide_content.content_blocks):
        # Calculate space needed for all remaining blocks after this one
        remaining_height = sum(
            block_estimated_heights[j] + gap
            for j in range(idx + 1, len(slide_content.content_blocks))
            if block_estimated_heights[j] > 0
        )
        
        if block.block_type == "bullets" and block.bullets:
            bottom = add_bullets_to_slide(
                slide, block.bullets,
                top=current_top,
                use_template_position=is_template,
                body_format=body_format
            )
            current_top = bottom + gap
        elif block.block_type == "table" and block.table:
            bottom = add_table_to_slide(
                slide, prs, block.table,
                top=current_top,
                use_template_position=is_template,
                bottom_reserve=remaining_height,
                body_format=body_format
            )
            current_top = bottom + gap
    
    return slide


def _add_bullet_formatting(paragraph, level: int, is_ordered: bool = False,
                           is_label: bool = False):
    """Add bullet point formatting to a textbox paragraph.
    
    Textboxes don't inherit bullet formatting from slide masters,
    so bullet XML (buFont, buChar, marL, indent) must be added
    explicitly for bullet characters and indentation to render.
    
    IMPORTANT: Call this AFTER all text runs and paragraph properties
    (e.g. space_after) have been set, because p.text= can rebuild
    <a:pPr> and discard earlier children.
    
    Args:
        paragraph: The python-pptx paragraph object to format.
        level: Nesting level (0 = top-level, 1 = nested, etc.).
        is_ordered: If True, suppress bullet char (number is already in text).
        is_label: If True, suppress bullet char and use minimal indent
                  (for section labels / non-bulleted paragraphs).
    """
    from lxml import etree
    
    pPr = paragraph._p.get_or_add_pPr()
    
    # Remove any existing bullet elements to avoid duplicates
    _bu_tags = [qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum'),
                qn('a:buFont'), qn('a:buFontTx'),
                qn('a:buSzPct'), qn('a:buSzPts'), qn('a:buSzTx'),
                qn('a:buClr')]
    for child in list(pPr):
        if child.tag in _bu_tags:
            pPr.remove(child)
    
    if is_label:
        # Label / non-bulleted paragraph: suppress bullet, minimal indent
        pPr.set('marL', '228600')  # ~0.25 inches left margin
        pPr.set('indent', '0')
        etree.SubElement(pPr, qn('a:buNone'))
        return
    
    # Indentation: margin left increases per nesting level, with a
    # negative first-line (hanging) indent so the bullet/number hangs left.
    #
    # Unordered bullets use a tighter margin (0.25") matching the template
    # and label paragraphs, so bullet text aligns with numbered-list text
    # and non-bulleted labels at the same horizontal position.
    # Ordered (numbered) items keep a wider margin (0.375") because the
    # inline "1. " prefix naturally offsets the content text.
    level_step = 457200     # 0.5 inches per nesting level in EMU
    
    if is_ordered:
        base_margin = 342900    # 0.375 inches in EMU
        hanging = -342900       # hanging indent in EMU
    else:
        base_margin = 228600    # 0.25 inches in EMU (matches labels / template)
        hanging = -228600       # hanging indent in EMU
    
    margin_left = base_margin + (level * level_step)
    pPr.set('marL', str(margin_left))
    pPr.set('indent', str(hanging))
    
    if is_ordered:
        # Numbered items already have the number prepended as a text run.
        # Add buNone to suppress any inherited bullet character.
        etree.SubElement(pPr, qn('a:buNone'))
    else:
        # Unordered bullets — add bullet font + character per nesting level.
        # buFont is required for PowerPoint to render the character in a textbox.
        bullet_chars = ['\u2022', '\u2013', '\u25aa']  # •, –, ▪
        char = bullet_chars[min(level, len(bullet_chars) - 1)]
        
        # Bullet color — Hawaii Makai (Ocean) theme for consistent bullet dots
        # buClr must come BEFORE buFont/buChar in the XML per OOXML schema
        bu_color = get_color(CONFIG["color_mapping"]["bullet_color"])
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % (bu_color[0], bu_color[1], bu_color[2]))
        
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', 'Arial')
        buFont.set('pitchFamily', '34')
        buFont.set('charset', '0')
        
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', char)


def add_bullets_to_slide(slide, bullets: List[BulletItem],
                         top: Optional[float] = None,
                         max_height: Optional[float] = None,
                         use_template_position: bool = False,
                         body_format: Optional[Dict[str, Any]] = None) -> float:
    """
    Add bullet points to a slide via a manual textbox.
    
    When *body_format* is provided (extracted from the template's body
    placeholder), those values are used for font name/size/color so the
    textbox bullets visually match the template.  Otherwise CONFIG defaults
    are used.
    
    Args:
        slide: The slide to add bullets to
        bullets: List of BulletItem objects
        top: Top position in inches (overrides use_template_position if set)
        max_height: Maximum height in inches for the text box
        use_template_position: If True, position content lower to avoid template title area
        body_format: Optional dict with 'font_name', 'font_size', 'font_color'
                     extracted from the template body placeholder
    
    Returns:
        Bottom Y position in inches after the bullet content
    """
    # Determine top position
    if top is not None:
        top_inches = top
    elif use_template_position:
        top_inches = 1.8
    else:
        top_inches = 1.5
    
    # Estimate bullet block height, accounting for text wrapping
    body_size_pt = CONFIG["fonts"]["body"]["size"]
    line_height = body_size_pt * 1.5 / 72  # single line height in inches
    space_after = 8 / 72  # 8pt space after each bullet, in inches
    
    # Available text width for estimating wraps (textbox width minus indent margins)
    text_box_width_inches = 12.333
    # Approximate characters per line at this font size
    # Open Sans at 18pt: ~0.11 inches per character average
    avg_char_width_inches = body_size_pt * 0.006  # empirical approximation
    chars_per_line = text_box_width_inches / avg_char_width_inches
    
    estimated_height = 0.0
    for bullet in bullets:
        # Total text length across all runs
        total_chars = sum(len(run.text) for run in bullet.text_runs)
        # Reduce available width for nested bullets (indented)
        indent_reduction = bullet.level * 0.5  # ~0.5 inches per indent level
        effective_chars_per_line = max(20, (text_box_width_inches - indent_reduction) / avg_char_width_inches)
        # Estimate number of wrapped lines (minimum 1)
        num_lines = max(1, -(-total_chars // int(effective_chars_per_line)))  # ceiling division
        estimated_height += (num_lines * line_height) + space_after
    
    estimated_height = max(0.5, estimated_height)
    
    # Apply max_height constraint
    if max_height is not None:
        box_height = min(estimated_height, max_height)
    else:
        box_height = min(estimated_height, 5.2)
    
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_inches), 
        Inches(12.333), Inches(box_height)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            # Clear any default content from the empty first paragraph
            # (preserves <a:pPr> if it exists)
            for child in list(p._p):
                if child.tag != qn('a:pPr'):
                    p._p.remove(child)
        else:
            p = tf.add_paragraph()
        
        # Add text runs via add_run() — avoids p.text= which can
        # rebuild <a:pPr> and discard bullet formatting
        for run in bullet.text_runs:
            r = p.add_run()
            r.text = run.text
            
            # Use template body formatting when available, else CONFIG defaults
            r.font.name = (body_format.get("font_name")
                           if body_format and body_format.get("font_name")
                           else CONFIG["fonts"]["body"]["name"])
            r.font.size = (body_format["font_size"]
                           if body_format and body_format.get("font_size")
                           else Pt(CONFIG["fonts"]["body"]["size"]))
            r.font.color.rgb = (body_format["font_color"]
                                if body_format and body_format.get("font_color")
                                else get_color(CONFIG["color_mapping"]["body_text"]))
            
            if run.bold:
                r.font.bold = True
            if run.italic:
                r.font.italic = True
            if run.link:
                r.hyperlink.address = run.link
        
        p.space_after = Pt(8)
        
        # Labels get extra space_before for visual separation from
        # preceding content (but not for the first item on the slide)
        if bullet.is_label and i > 0:
            p.space_before = Pt(12)
        
        # Add bullet formatting LAST — after all text runs and paragraph
        # properties are set, so nothing can overwrite the bullet XML
        _add_bullet_formatting(p, bullet.level, bullet.is_ordered, bullet.is_label)
    
    return top_inches + estimated_height


def parse_cell_formatting(cell_text: str) -> List[tuple]:
    """
    Parse markdown formatting in table cell text.
    Returns list of (text, is_bold, is_italic) tuples.
    """
    result = []
    text = str(cell_text)
    current_pos = 0
    
    # Pattern for **bold** and *italic*
    pattern = r'(\*\*.*?\*\*|\*.*?\*)'
    
    for match in re.finditer(pattern, text):
        # Add text before match
        if match.start() > current_pos:
            result.append((text[current_pos:match.start()], False, False))
        
        matched_text = match.group(1)
        
        # Bold
        if matched_text.startswith('**') and matched_text.endswith('**'):
            result.append((matched_text[2:-2], True, False))
        # Italic
        elif matched_text.startswith('*') and matched_text.endswith('*'):
            result.append((matched_text[1:-1], False, True))
        
        current_pos = match.end()
    
    # Add remaining text
    if current_pos < len(text):
        result.append((text[current_pos:], False, False))
    
    return result if result else [(text, False, False)]


def _set_cell_border(cell, side: str, width_pt: float, color: RGBColor):
    """Set a single border on a table cell via XML manipulation.
    
    Args:
        cell: python-pptx table cell
        side: 'L', 'R', 'T', or 'B' (left, right, top, bottom)
        width_pt: Border width in points (e.g. 0.5, 1.0, 2.0)
        color: RGBColor for the border line
    """
    from lxml import etree
    
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    
    tag = qn(f'a:ln{side}')
    # Remove existing border element for this side
    for existing in tcPr.findall(tag):
        tcPr.remove(existing)
    
    ln = etree.SubElement(tcPr, tag)
    ln.set('w', str(int(width_pt * 12700)))  # points to EMU
    ln.set('cmpd', 'sng')  # single line
    
    solidFill = etree.SubElement(ln, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', str(color))


def _clear_cell_border(cell, side: str):
    """Remove a border from a table cell (set to no line).
    
    Args:
        cell: python-pptx table cell
        side: 'L', 'R', 'T', or 'B'
    """
    from lxml import etree
    
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    
    tag = qn(f'a:ln{side}')
    for existing in tcPr.findall(tag):
        tcPr.remove(existing)
    
    ln = etree.SubElement(tcPr, tag)
    ln.set('w', '0')
    etree.SubElement(ln, qn('a:noFill'))


def _apply_modern_table_borders(pptx_table, num_rows: int, num_cols: int,
                                 header_border_color: RGBColor,
                                 row_border_color: RGBColor):
    """Apply modern clean-line table borders.
    
    Style: thick dark border under header row, thin subtle horizontal borders
    between data rows, no vertical borders.  This gives a clean, professional
    consulting-style table appearance.
    
    Args:
        pptx_table: The python-pptx Table object
        num_rows: Total number of rows (including header)
        num_cols: Number of columns
        header_border_color: Color for the header bottom border
        row_border_color: Color for data row borders
    """
    for row_idx in range(num_rows):
        for col_idx in range(num_cols):
            cell = pptx_table.cell(row_idx, col_idx)
            
            # Clear all vertical borders (left and right)
            _clear_cell_border(cell, 'L')
            _clear_cell_border(cell, 'R')
            
            if row_idx == 0:
                # Header row: clear top, add thick bottom border
                _clear_cell_border(cell, 'T')
                _set_cell_border(cell, 'B', 1.5, header_border_color)
            elif row_idx == num_rows - 1:
                # Last row: thin top border, thin bottom border
                _set_cell_border(cell, 'T', 0.5, row_border_color)
                _set_cell_border(cell, 'B', 0.75, row_border_color)
            else:
                # Middle data rows: thin top and bottom borders
                _set_cell_border(cell, 'T', 0.5, row_border_color)
                _set_cell_border(cell, 'B', 0.5, row_border_color)


def set_cell_text_with_formatting(cell, text: str, font_size: int, font_name: str, font_color):
    """
    Set cell text with markdown formatting support (bold/italic).
    """
    parts = parse_cell_formatting(text)
    
    tf = cell.text_frame
    p = tf.paragraphs[0]
    
    for i, (part_text, is_bold, is_italic) in enumerate(parts):
        if i == 0:
            p.text = part_text
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = font_color
                run.font.bold = is_bold
                run.font.italic = is_italic
        else:
            run = p.add_run()
            run.text = part_text
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = font_color
            run.font.bold = is_bold
            run.font.italic = is_italic


def _detect_uniform_columns(table: TableData) -> set:
    """Detect columns that contain similar short data and should share equal widths.

    When a table has multiple data columns with similarly short numeric content
    (e.g. all percentage values, all small counts), those columns look best
    with equal widths and center alignment rather than proportional sizing.

    Detection criteria — a column is a **uniform candidate** if:
      • ≥60 % of its body cells are numeric (digits, commas, %, $, ~)
      • The longest cell text is ≤ 12 characters

    A uniform **group** forms when ≥ 2 candidates exist and their max cell
    lengths are within a 3× ratio of each other.

    Returns:
        Set of column indices that should be equalised / center-aligned.
        Empty set if no uniform group is detected.
    """
    num_cols = len(table.headers)
    if num_cols < 3:
        # Need at least 1 label column + 2 data columns
        return set()

    _num_re = re.compile(
        r'^[~$±]?[\d,]+\.?\d*\s*[%]?$'   # "12,345", "$99", "~5%", "3.14"
        r'|^\d+\s*[%]$'                    # "77%"
        r'|^[0-9]+$'                        # plain integer
    )

    candidates = []
    for col_idx in range(num_cols):
        max_len = 0
        numeric_count = 0
        for row in table.rows:
            if col_idx < len(row):
                text = row[col_idx].replace('**', '').replace('*', '').strip()
                max_len = max(max_len, len(text))
                if _num_re.match(text):
                    numeric_count += 1

        is_numeric = len(table.rows) > 0 and numeric_count / len(table.rows) >= 0.6
        is_short = max_len <= 12

        if is_numeric and is_short:
            candidates.append((col_idx, max_len))

    if len(candidates) < 2:
        return set()

    # Verify similarity: max-to-min content-length ratio ≤ 3
    max_lengths = [ml for _, ml in candidates]
    min_l = min(max_lengths) if max_lengths else 0
    max_l = max(max_lengths) if max_lengths else 0
    if min_l > 0 and max_l / min_l > 3.0:
        return set()

    return {idx for idx, _ in candidates}


def _calculate_column_widths(table: TableData, total_width_emu: int,
                             header_font_size: int, body_font_size: int,
                             uniform_cols: Optional[set] = None) -> list:
    """Calculate smart column widths that minimise wrapping in short columns.

    Strategy — "fit short columns first":
      1. Compute the natural (single-line, no-wrap) width each column needs.
      2. If every column fits on one line within the available width, scale up
         proportionally — all content renders without wrapping.
      3. Otherwise, *lock* columns whose natural width is small enough that
         they can display their content without wrapping, then distribute
         the remaining space among the wider columns proportionally.

    This prevents label-style columns (e.g. "Commodity") from being starved
    to a sliver while description columns get excess space.

    When *uniform_cols* is provided, those columns are equalised to the
    widest among them (e.g. multiple percentage columns).

    Args:
        table: TableData with headers and rows
        total_width_emu: Total table width in EMU
        header_font_size: Header font size in points
        body_font_size: Body font size in points
        uniform_cols: Optional set of column indices to equalise

    Returns:
        List of column widths in EMU
    """
    num_cols = len(table.headers)
    if num_cols == 0:
        return []

    # Character width in EMU: ~52% of font em for proportional sans-serif
    char_w_emu = int(body_font_size * 0.52 * 12700)  # 1 pt = 12700 EMU
    cell_margins_emu = 274320  # ~0.30" combined L+R cell margins

    # ---- Step 1: natural (single-line) width per column ----
    natural_widths = []
    for col_idx in range(num_cols):
        header_text = table.headers[col_idx].replace('**', '').replace('*', '')
        max_len = len(header_text)
        for row in table.rows:
            if col_idx < len(row):
                cell_text = row[col_idx].replace('**', '').replace('*', '')
                max_len = max(max_len, len(cell_text))
        natural_w = max(1, max_len) * char_w_emu + cell_margins_emu
        natural_widths.append(natural_w)

    total_natural = sum(natural_widths)

    if total_natural <= total_width_emu:
        # ---- All content fits on one line — scale up proportionally ----
        scale = total_width_emu / total_natural
        widths = [int(w * scale) for w in natural_widths]
    else:
        # ---- Some columns must wrap — lock short columns first ----
        # Threshold: lock a column if its natural width is ≤ this fraction
        # of total table width (generous for few columns, tighter for many)
        if num_cols <= 2:
            max_lock_frac = 0.35
        elif num_cols <= 4:
            max_lock_frac = 0.28
        else:
            max_lock_frac = 0.20
        max_lock_width = int(total_width_emu * max_lock_frac)
        min_col_width = int(total_width_emu * 0.08)  # absolute floor

        widths = [0] * num_cols
        locked = set()
        remaining_width = total_width_emu

        # Lock smallest columns first (greedy): ensures label-like columns
        # get their full natural width without wrapping
        for idx in sorted(range(num_cols), key=lambda i: natural_widths[i]):
            nw = natural_widths[idx]
            unlocked_remaining = num_cols - len(locked) - 1

            if nw <= max_lock_width:
                # Only lock if remaining space still gives other columns
                # at least min_col_width each
                leftover = remaining_width - nw
                if unlocked_remaining == 0 or leftover >= unlocked_remaining * min_col_width:
                    widths[idx] = nw
                    locked.add(idx)
                    remaining_width -= nw

        # Distribute remaining width to unlocked columns proportionally
        # to their natural widths (wider content gets more space)
        unlocked = [i for i in range(num_cols) if i not in locked]
        if unlocked:
            unlocked_natural = sum(natural_widths[i] for i in unlocked)
            if unlocked_natural > 0:
                for i in unlocked:
                    share = remaining_width * natural_widths[i] / unlocked_natural
                    widths[i] = max(min_col_width, int(share))
            else:
                per_col = remaining_width // len(unlocked)
                for i in unlocked:
                    widths[i] = per_col

    # ---- Enforce header minimum widths ----
    header_min_widths = []
    header_char_w_emu = int(header_font_size * 0.52 * 12700)
    for col_idx in range(num_cols):
        header_text = table.headers[col_idx].replace('**', '').replace('*', '')
        min_w = int(len(header_text) * header_char_w_emu + cell_margins_emu)
        header_min_widths.append(min_w)

    for i in range(num_cols):
        deficit = header_min_widths[i] - widths[i]
        if deficit > 0:
            widest_idx = widths.index(max(widths))
            if widest_idx != i and widths[widest_idx] - deficit > header_min_widths[widest_idx]:
                widths[i] += deficit
                widths[widest_idx] -= deficit

    # ---- Equalise uniform column groups ----
    if uniform_cols and len(uniform_cols) >= 2:
        u_indices = sorted(uniform_cols)
        non_u = [i for i in range(num_cols) if i not in uniform_cols]

        # Floor widths for non-uniform columns: use their natural (single-line)
        # width so uniform equalization never compresses label columns into
        # wrapping.  natural_widths[i] already accounts for the longest cell.
        non_u_floors = {}
        for i in non_u:
            non_u_floors[i] = max(header_min_widths[i], natural_widths[i])

        target_w = max(widths[i] for i in u_indices)
        extra_needed = sum(max(0, target_w - widths[i]) for i in u_indices)

        if extra_needed > 0 and non_u:
            total_non_u = sum(widths[i] for i in non_u)
            total_floor = sum(non_u_floors[i] for i in non_u)
            max_shrink = max(0, total_non_u - total_floor)
            if extra_needed > max_shrink:
                available_for_uniform = total_width_emu - total_floor
                target_w = available_for_uniform // len(u_indices)
                extra_needed = sum(max(0, target_w - widths[i]) for i in u_indices)

        for i in u_indices:
            widths[i] = target_w

        if extra_needed > 0 and non_u:
            total_non_u = sum(widths[i] for i in non_u)
            if total_non_u > 0:
                for i in non_u:
                    shrink = int(extra_needed * widths[i] / total_non_u)
                    widths[i] = max(non_u_floors[i], widths[i] - shrink)

    # Distribute rounding remainder to the widest column
    remainder = total_width_emu - sum(widths)
    if remainder != 0:
        widest_idx = widths.index(max(widths))
        widths[widest_idx] += remainder

    return widths


def _estimate_natural_table_width(table: TableData, header_font_size: int,
                                   body_font_size: int,
                                   cell_margin_left_emu: int,
                                   cell_margin_right_emu: int) -> int:
    """Estimate the natural width a table needs based on its content.

    Calculates per-column width from the longest cell text (header or body),
    using approximate character-width metrics for proportional sans-serif
    fonts.  The sum of all columns gives the natural (unwrapped) table width
    in EMU.

    Args:
        table: TableData with headers and rows
        header_font_size: Header font size in points (after scaling)
        body_font_size: Body font size in points (after scaling)
        cell_margin_left_emu: Left cell margin in EMU
        cell_margin_right_emu: Right cell margin in EMU

    Returns:
        Estimated natural table width in EMU
    """
    num_cols = len(table.headers)
    if num_cols == 0:
        return 0

    # Approximate character width: ~52% of font size for proportional
    # sans-serif fonts (Calibri, Arial).  1 pt = 12700 EMU.
    header_char_width = int(header_font_size * 0.52 * 12700)
    body_char_width = int(body_font_size * 0.52 * 12700)

    cell_h_margins = cell_margin_left_emu + cell_margin_right_emu

    total_natural = 0
    for col_idx in range(num_cols):
        # Header text width (strip markdown bold/italic markers)
        header_text = table.headers[col_idx].replace('**', '').replace('*', '')
        header_width = len(header_text) * header_char_width

        # Maximum body cell width in this column
        max_body_width = 0
        for row in table.rows:
            if col_idx < len(row):
                cell_text = row[col_idx].replace('**', '').replace('*', '')
                cell_width = len(cell_text) * body_char_width
                max_body_width = max(max_body_width, cell_width)

        # Natural column width = max(header, body) + cell margins
        col_natural = max(header_width, max_body_width) + cell_h_margins
        total_natural += col_natural

    return total_natural


def add_table_to_slide(slide, prs: Presentation, table: TableData,
                       top: Optional[float] = None,
                       use_template_position: bool = False,
                       bottom_reserve: float = 0.0,
                       body_format: Optional[Dict[str, Any]] = None) -> float:
    """
    Add a table to a slide with proper sizing and formatting.
    
    Args:
        slide: The slide to add the table to
        prs: Presentation object for dimensions
        table: TableData object with headers and rows
        top: Top position in inches (overrides use_template_position if set)
        use_template_position: If True, position content lower to avoid template title area
        bottom_reserve: Extra inches to reserve below the table for subsequent
                        content blocks (e.g. bullet points after the table).
        body_format: Optional dict with 'font_color' extracted from the template
                     body placeholder, used to keep table text color consistent
                     with bullet text color.
    
    Returns:
        Bottom Y position in inches after the table
    """
    num_cols = len(table.headers)
    num_rows = len(table.rows) + 1
    
    # Calculate table dimensions to fit within slide
    left_margin = Inches(0.5)
    right_margin = Inches(0.8)  # Extra space for page number
    
    # Template slides need more bottom margin to clear footer/accent elements
    bottom_margin = Inches(0.9) if use_template_position else Inches(0.4)
    
    # Available width (slide width minus margins)
    slide_width_inches = prs.slide_width.inches
    table_width = Inches(slide_width_inches - left_margin.inches - right_margin.inches)
    
    # Adjust position based on parameters
    if top is not None:
        top_inches = top
    elif use_template_position:
        top_inches = 1.5
    else:
        top_inches = 1.4
    top_position = Inches(top_inches)
    
    # Available height (slide height minus top position, bottom margin, and reserved space)
    slide_height_inches = prs.slide_height.inches
    total_below_top = max(0.5, slide_height_inches - top_inches - bottom_margin.inches)
    
    # Smart bottom_reserve cap: ensure the table gets at least a minimum
    # fraction of available space.  Without this, over-estimated bullet
    # heights can starve small tables, producing tiny text and cramped rows.
    if num_rows <= 5:
        min_table_fraction = 0.50   # small tables get ≥50% of space
    elif num_rows <= 8:
        min_table_fraction = 0.55
    else:
        min_table_fraction = 0.65   # large tables get ≥65% of space
    
    min_table_space = total_below_top * min_table_fraction
    effective_reserve = min(bottom_reserve, total_below_top - min_table_space)
    effective_reserve = max(0.0, effective_reserve)
    available_height = max(0.5, total_below_top - effective_reserve)
    
    if effective_reserve < bottom_reserve:
        logger.debug("    Table reserve capped: requested=%.2f, effective=%.2f (table gets %.1f%% of %.2f\")",
                     bottom_reserve, effective_reserve, min_table_fraction * 100, total_below_top)
    
    # Calculate row height based on number of rows
    # Use tighter minimum for large tables to ensure they fit
    min_row_height = 0.28 if num_rows > 8 else 0.35  # inches
    # Dynamic max — small tables can have taller, more spacious rows
    if num_rows <= 5:
        max_row_height = 0.80
    elif num_rows <= 8:
        max_row_height = 0.70
    else:
        max_row_height = 0.55
    calculated_row_height = available_height / num_rows
    row_height = max(min_row_height, min(max_row_height, calculated_row_height))
    
    # Clamp table height to available space
    actual_table_height = min(available_height, row_height * num_rows)
    table_height = Inches(actual_table_height)
    
    # ---- Font sizes & cell margins (computed before table creation so we ----
    # ---- can estimate natural table width from content)                  ----
    
    # Target: match body text size; reduce only when table content won't fit
    target_font_size = CONFIG["fonts"]["body"]["size"]
    is_large_table = num_rows > 8
    
    # Cell margin settings — generous left/right margins for readability
    cell_margin_top = Inches(0.03) if is_large_table else Inches(0.06)
    cell_margin_bottom = Inches(0.03) if is_large_table else Inches(0.06)
    cell_margin_left = Inches(0.10) if is_large_table else Inches(0.15)
    cell_margin_right = Inches(0.06) if is_large_table else Inches(0.10)
    
    # --- Intelligent font sizing: start at body size, reduce only if needed ---
    # Estimate proportional column widths from max content length per column
    table_width_inches = int(table_width) / 914400
    col_max_chars = []
    for col_idx in range(num_cols):
        max_len = len(table.headers[col_idx].replace('**', '').replace('*', ''))
        for row in table.rows:
            if col_idx < len(row):
                cell_text = row[col_idx].replace('**', '').replace('*', '')
                max_len = max(max_len, len(cell_text))
        col_max_chars.append(max(1, max_len))
    total_chars = sum(col_max_chars)
    h_margin_per_col = cell_margin_left.inches + cell_margin_right.inches
    usable_width = max(1.0, table_width_inches - num_cols * h_margin_per_col)
    col_widths_est = [max(0.5, (c / total_chars) * usable_width) for c in col_max_chars]
    
    # Vertical overhead per row (cell margins + breathing room)
    vert_overhead = cell_margin_top.inches + cell_margin_bottom.inches + 0.03
    
    # Step down from target font size until estimated table height fits
    # in the available vertical space.  Line height ≈ 1.2× font size for
    # single-spaced Calibri; char width ≈ 0.52× font size.
    table_font_size = target_font_size
    est_total_height = 0.0
    while table_font_size > 9:
        line_ht = table_font_size * 1.2 / 72      # single-line height (inches)
        char_w = table_font_size * 0.52 / 72       # approx char width (inches)
        
        # Estimate total table height at this font size (header + data rows)
        est_total_height = 0.0
        for row_idx in range(-1, len(table.rows)):   # -1 = header row
            cells = table.headers if row_idx == -1 else table.rows[row_idx]
            row_max_lines = 1
            for col_idx in range(min(num_cols, len(cells))):
                cell_text = cells[col_idx].replace('**', '').replace('*', '')
                if not cell_text:
                    continue
                col_w = col_widths_est[col_idx] if col_idx < len(col_widths_est) else 2.0
                cpl = max(1, int(col_w / char_w))            # chars per line
                lines = max(1, -(-len(cell_text) // cpl))    # ceil division
                row_max_lines = max(row_max_lines, lines)
            est_total_height += row_max_lines * line_ht + vert_overhead
        
        if est_total_height <= available_height:
            break
        table_font_size -= 1
    
    header_font_size = table_font_size
    body_font_size = table_font_size
    logger.debug("    Font scaling: target=%dpt -> table=%dpt (est_h=%.2f\", avail=%.2f\")",
                 target_font_size, table_font_size, est_total_height, available_height)
    
    # ---- Content-aware table width: shrink to fit content, center ----
    available_width_emu = int(table_width)
    natural_width_emu = _estimate_natural_table_width(
        table, header_font_size, body_font_size,
        int(cell_margin_left), int(cell_margin_right)
    )
    # Add breathing room (20%) so text isn't jammed edge-to-edge
    padded_natural_emu = int(natural_width_emu * 1.20)
    
    # Clamp: at least 50% of available width, at most 100%
    min_table_width_emu = int(available_width_emu * 0.50)
    actual_table_width_emu = max(min_table_width_emu,
                                 min(padded_natural_emu, available_width_emu))
    
    # Center the table horizontally if narrower than full width
    if actual_table_width_emu < available_width_emu:
        centering_offset_emu = (available_width_emu - actual_table_width_emu) // 2
        left_position = Emu(int(left_margin) + centering_offset_emu)
        actual_table_width = Emu(actual_table_width_emu)
        logger.debug("    Table width: natural=%.1f\", padded=%.1f\", "
                     "actual=%.1f\" (%.0f%% of available), centered",
                     natural_width_emu / 914400, padded_natural_emu / 914400,
                     actual_table_width_emu / 914400,
                     actual_table_width_emu / available_width_emu * 100)
    else:
        left_position = left_margin
        actual_table_width = table_width
    
    pptx_table = slide.shapes.add_table(
        num_rows, num_cols, 
        left_position, top_position, 
        actual_table_width, table_height
    ).table
    
    # Detect uniform column groups (e.g. several short % columns) that
    # should share equal widths and be center-aligned.
    uniform_cols = _detect_uniform_columns(table)

    # Content-aware column widths (proportional to cell content length,
    # with uniform groups equalised)
    col_widths = _calculate_column_widths(table, int(actual_table_width),
                                          header_font_size, body_font_size,
                                          uniform_cols=uniform_cols)
    for i in range(num_cols):
        pptx_table.columns[i].width = col_widths[i]
    
    # Detect numeric columns for right-alignment.
    # A column is "numeric" if the majority of its data cells look like numbers
    # (digits, commas, periods, %, $, ~, +/-, or short words like "Yes"/"No").
    _numeric_pattern = re.compile(r'^[~$]?[\d,]+\.?\d*\s*[%]?$|^\d+[%]$')
    numeric_cols = set()
    for col_idx in range(num_cols):
        numeric_count = 0
        for row in table.rows:
            if col_idx < len(row):
                cell_text = row[col_idx].replace('**', '').replace('*', '').strip()
                if _numeric_pattern.match(cell_text):
                    numeric_count += 1
        if len(table.rows) > 0 and numeric_count / len(table.rows) >= 0.6:
            numeric_cols.add(col_idx)
    
    # Header row
    for i, header in enumerate(table.headers):
        cell = pptx_table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["table_header_bg"])
        
        # Cell margins
        cell.margin_top = cell_margin_top
        cell.margin_bottom = cell_margin_bottom
        cell.margin_left = cell_margin_left
        cell.margin_right = cell_margin_right
        
        p = cell.text_frame.paragraphs[0]
        p.font.bold = CONFIG["fonts"]["table_header"]["bold"]
        p.font.size = Pt(header_font_size)
        p.font.name = CONFIG["fonts"]["table_header"]["name"]
        p.font.color.rgb = get_color(CONFIG["color_mapping"]["table_header_text"])
        # Alignment: uniform group → CENTER, other numeric → RIGHT, else LEFT
        if i in uniform_cols:
            p.alignment = PP_ALIGN.CENTER
        elif i in numeric_cols:
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        # Vertical alignment and word wrap
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text_frame.word_wrap = True
    
    # Data rows with formatting support
    for row_idx, row in enumerate(table.rows):
        for col_idx in range(num_cols):
            cell = pptx_table.cell(row_idx + 1, col_idx)
            cell_text = row[col_idx] if col_idx < len(row) else ""
            
            # Cell margins
            cell.margin_top = cell_margin_top
            cell.margin_bottom = cell_margin_bottom
            cell.margin_left = cell_margin_left
            cell.margin_right = cell_margin_right
            
            # Use formatted text setting to handle **bold** markers
            # Use template body color when available for consistency with bullets
            table_text_color = (body_format["font_color"]
                                if body_format and body_format.get("font_color")
                                else get_color(CONFIG["color_mapping"]["body_text"]))
            set_cell_text_with_formatting(
                cell, 
                cell_text,
                body_font_size,
                CONFIG["fonts"]["table_body"]["name"],
                table_text_color
            )
            
            # Hanging indent for numbered items (e.g. "1. Text" or "10. Text")
            # so wrapped lines align with the text after the number prefix.
            plain_text = cell_text.replace('**', '').replace('*', '')
            num_match = re.match(r'^(\d+\.\s)', plain_text)
            if num_match:
                prefix = num_match.group(1)  # e.g. "4. "
                # Calculate prefix width by character type (digits are wider,
                # period and space are narrow).  Units: fraction of font em.
                prefix_width_em = 0.0
                for ch in prefix:
                    if ch.isdigit():
                        prefix_width_em += 0.60   # digit width ~0.6 em
                    elif ch == '.':
                        prefix_width_em += 0.28   # period ~0.28 em
                    elif ch == ' ':
                        prefix_width_em += 0.28   # space ~0.28 em
                    else:
                        prefix_width_em += 0.55   # fallback
                # Convert em-widths to EMU (1 pt = 12700 EMU)
                hanging_emu = int(prefix_width_em * body_font_size * 12700)
                for para in cell.text_frame.paragraphs:
                    pPr = para._p.get_or_add_pPr()
                    pPr.set('marL', str(hanging_emu))
                    pPr.set('indent', str(-hanging_emu))
            
            # Vertical alignment, word wrap, and text alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            if col_idx in uniform_cols:
                align = PP_ALIGN.CENTER
            elif col_idx in numeric_cols:
                align = PP_ALIGN.RIGHT
            else:
                align = PP_ALIGN.LEFT
            for para in cell.text_frame.paragraphs:
                para.alignment = align
            
            # Alternate row colors (even rows get light gray fill,
            # odd rows get a subtle off-white so they don't blend
            # into the white slide background)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["table_alt_row"])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 244, 247)  # off-white #F2F4F7
    
    # Disable PowerPoint's built-in banding so our custom styling prevails
    pptx_table.first_row = False
    pptx_table.horz_banding = False
    pptx_table.vert_banding = False
    pptx_table.last_row = False
    pptx_table.first_col = False
    pptx_table.last_col = False
    
    # Apply modern clean-line borders
    _apply_modern_table_borders(
        pptx_table, num_rows, num_cols,
        header_border_color=get_color(CONFIG["color_mapping"]["table_header_bg"]),
        row_border_color=RGBColor(200, 200, 200)  # subtle light gray
    )
    
    return top_inches + actual_table_height


# ============================================
# MAIN CONVERSION FUNCTION
# ============================================

def convert_markdown_to_pptx(
    markdown_path: str, 
    output_path: Optional[str] = None,
    template_path: Optional[str] = None
) -> str:
    """
    Convert markdown file to PowerPoint presentation
    
    Args:
        markdown_path: Path to input markdown file
        output_path: Path to output PPTX file (optional)
        template_path: Path to PowerPoint template file (optional)
    
    Returns:
        Path to generated PPTX file
    """
    # Read markdown file
    with open(markdown_path, 'r', encoding='utf-8') as f:
        markdown_text = f.read()
    
    # Parse markdown
    logger.info("Parsing markdown...")
    parser = MarkdownToSlides(markdown_text)
    slides = parser.parse()
    logger.info("Found %d slides", len(slides))
    
    # Create presentation with layout manager
    logger.info("Creating presentation...")
    prs, layout_manager = create_presentation(template_path)
    
    if layout_manager:
        logger.info("Intelligent layout selection enabled with %d layouts available", len(layout_manager.layouts))
    else:
        logger.info("Using manual slide creation (no template)")
    
    # Generate slides with intelligent layout selection
    logger.info("Generating slides...")
    for i, slide_content in enumerate(slides):
        title_display = (slide_content.title[:50] + "...") if len(slide_content.title) > 50 else slide_content.title
        logger.info("Slide %d: %s - '%s'", i + 1, slide_content.slide_type.value, title_display)
        
        if slide_content.slide_type == SlideType.TITLE:
            add_title_slide(prs, slide_content.title, slide_content.subtitle, layout_manager)
        elif slide_content.slide_type == SlideType.SECTION:
            add_section_slide(prs, slide_content.title, layout_manager)
        elif slide_content.slide_type == SlideType.CONTENT:
            add_content_slide(prs, slide_content, layout_manager)
    
    # Ensure slide-number placeholders are present on every slide
    # (python-pptx doesn't copy decorator placeholders from the layout)
    if layout_manager:
        for slide in prs.slides:
            _ensure_slide_number(slide)

    # Determine output path
    if not output_path:
        base = os.path.splitext(markdown_path)[0]
        output_path = f"{base}.pptx"
    
    # Clear metadata and save presentation
    _clear_metadata(prs)
    prs.save(output_path)
    logger.info("=" * 50)
    logger.info("Presentation saved to: %s", output_path)
    logger.info("Total slides: %d", len(prs.slides))
    
    return output_path


# ============================================
# COMMAND LINE INTERFACE
# ============================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert markdown files to PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python md_to_pptx_converter.py document.md
  python md_to_pptx_converter.py document.md --template template.pptx
  python md_to_pptx_converter.py document.md -o output.pptx
  
  # Create a clean template (no slides) from an existing one:
  python md_to_pptx_converter.py --clean-template source.pptx -o clean_template.pptx
        """
    )
    
    parser.add_argument(
        'markdown_file',
        nargs='?',  # Make optional for --clean-template mode
        help='Path to input markdown file'
    )
    
    parser.add_argument(
        '-o', '--output',
        help='Path to output PPTX file (default: same name as input with .pptx extension)'
    )
    
    parser.add_argument(
        '-t', '--template',
        help='Path to PowerPoint template file (.pptx or .potx)'
    )
    
    parser.add_argument(
        '--clean-template',
        metavar='SOURCE',
        help='Create a clean template (no slides) from SOURCE template file'
    )
    
    parser.add_argument(
        '-v', '--verbose',
        action='store_true',
        help='Enable verbose/debug output'
    )
    
    args = parser.parse_args()
    
    # Configure logging based on verbosity
    log_level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(
        level=log_level,
        format='%(message)s'
    )
    
    # Handle --clean-template mode
    if args.clean_template:
        if not os.path.exists(args.clean_template):
            logger.error("Error: Source template not found: %s", args.clean_template)
            return 1
        
        output_path = args.output or args.clean_template.replace('.pptx', '_clean.pptx')
        try:
            create_clean_template(args.clean_template, output_path)
            return 0
        except Exception as e:
            logger.error("Error creating clean template: %s", e)
            return 1
    
    # Normal conversion mode requires markdown file
    if not args.markdown_file:
        parser.print_help()
        return 1
    
    # Validate input file
    if not os.path.exists(args.markdown_file):
        logger.error("Error: Input file not found: %s", args.markdown_file)
        return 1
    
    # Auto-discover template if not explicitly provided:
    # Look for .pptx/.potx files in the same directory as this script.
    if not args.template:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        for candidate in sorted(os.listdir(script_dir)):
            if candidate.lower().endswith(('.pptx', '.potx')) and 'template' in candidate.lower():
                args.template = os.path.join(script_dir, candidate)
                logger.info("Auto-detected template: %s", args.template)
                break
    
    # Validate template if provided
    if args.template and not os.path.exists(args.template):
        logger.warning("Warning: Template file not found: %s", args.template)
        logger.warning("Proceeding with default styling...")
        args.template = None
    
    # Convert
    try:
        output_path = convert_markdown_to_pptx(
            args.markdown_file,
            args.output,
            args.template
        )
        logger.info("Success! Created: %s", output_path)
        return 0
    except Exception as e:
        logger.error("Error: %s", e)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())
