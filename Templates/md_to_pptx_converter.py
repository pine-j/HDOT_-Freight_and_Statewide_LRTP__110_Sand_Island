"""
Markdown to PowerPoint Converter

Converts structured markdown files into professional PowerPoint presentations
with support for HDOT branding and custom templates.

Usage:
    python md_to_pptx_converter.py input.md
    python md_to_pptx_converter.py input.md --template template.potx
    python md_to_pptx_converter.py input.md -o output.pptx

Requirements:
    pip install python-pptx

Author: Generated for HDOT presentations
"""

import argparse
import logging
import os
import re
import shutil
import tempfile
import warnings
import zipfile
from dataclasses import dataclass, field
from enum import Enum
from typing import List, Optional, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor

# Configure module-level logger
logger = logging.getLogger(__name__)

# Suppress harmless duplicate name warnings when clearing template slides
warnings.filterwarnings('ignore', message='Duplicate name:.*ppt/slides/')

# ============================================
# CONFIGURATION - Modify these settings
# ============================================

CONFIG = {
    # Official HDOT Brand Colors (from State of Hawaii Web Style Guide)
    # These are ONLY used when template_path is None
    # Using Makai (Ocean) theme as primary - best for general HDOT branding
    "colors": {
        # Makai (Ocean) Theme - Primary
        "primary_blue": (38, 72, 110),      # #26486E - Dark Blue (primary brand color)
        "medium_blue": (77, 118, 164),      # #4D76A4 - Medium Blue (secondary)
        "dark_gray": (85, 85, 85),          # #555555 - Body text, neutral accents
        
        # Standard colors
        "white": (255, 255, 255),           # #ffffff
        "light_gray": (218, 222, 229),      # #dadee5 - Light gray for alternating rows
        
        # Mauka (Mountain) Theme - Alternative for inland projects
        "dark_green": (47, 89, 26),         # #2F591A
        "medium_green": (46, 114, 20),      # #2E7214
        "light_green": (128, 192, 98),      # #80C062
        
        # Aina (Land/Earth) Theme - Alternative for sustainability
        "dark_teal": (23, 84, 80),          # #175450
        "teal": (34, 118, 124),             # #22767C
        "light_teal": (115, 201, 207),      # #73C9CF
        
        # Sandy (Soil/Sand) Theme - Alternative for construction
        "dark_brown": (85, 65, 28),         # #55411C
        "medium_brown": (137, 113, 72),     # #897148
        "beige": (191, 183, 137),           # #BFB789
    },
    
    # Color assignments for slide elements (when NOT using template)
    "color_mapping": {
        "title_slide_bg": "primary_blue",      # Dark Blue #26486E
        "title_slide_text": "white",
        "section_slide_bar": "medium_blue",    # Medium Blue #4D76A4
        "section_slide_text": "white",
        "content_title_bar": "primary_blue",   # Dark Blue #26486E
        "content_title_text": "white",
        "body_text": "dark_gray",              # Dark Gray #555555
        "table_header_bg": "primary_blue",     # Dark Blue #26486E
        "table_header_text": "white",
        "table_alt_row": "light_gray",         # Light Gray #dadee5
        "accent": "medium_blue",               # Medium Blue #4D76A4
    },
    
    # Font settings (HDOT standard: Open Sans)
    "fonts": {
        "title": {"name": "Open Sans", "size": 44, "bold": True},
        "section": {"name": "Open Sans", "size": 36, "bold": True},
        "slide_title": {"name": "Open Sans", "size": 28, "bold": True},
        "body": {"name": "Open Sans", "size": 18},
        "table_header": {"name": "Open Sans", "size": 14, "bold": True},  # Increased from 12
        "table_body": {"name": "Open Sans", "size": 12},  # Increased from 11
    },
    
    # Slide dimensions (widescreen 16:9)
    "slide_width": 13.333,  # inches
    "slide_height": 7.5,    # inches
}


# ============================================
# DATA STRUCTURES
# ============================================

class SlideType(Enum):
    """Enumeration of supported slide types"""
    TITLE = "title"
    SECTION = "section"
    CONTENT = "content"


# Placeholder type constants for readability
PH_TITLE = PP_PLACEHOLDER.TITLE           # 1
PH_BODY = PP_PLACEHOLDER.BODY             # 2
PH_CENTER_TITLE = PP_PLACEHOLDER.CENTER_TITLE  # 3
PH_SUBTITLE = PP_PLACEHOLDER.SUBTITLE     # 4
PH_OBJECT = PP_PLACEHOLDER.OBJECT         # 7
PH_PICTURE = PP_PLACEHOLDER.PICTURE       # 18

# Media clip placeholder type (no enum constant in python-pptx)
PH_MEDIA_CLIP = 14


@dataclass
class TextRun:
    """Represents a formatted text run"""
    text: str
    bold: bool = False
    italic: bool = False
    link: Optional[str] = None


@dataclass
class BulletItem:
    """Represents a bullet point with optional nesting"""
    text_runs: List[TextRun]
    level: int = 0  # 0 = top level, 1 = nested, etc.


@dataclass
class TableData:
    """Represents a table"""
    headers: List[str]
    rows: List[List[str]]


@dataclass
class SlideContent:
    """Represents the content of a slide"""
    slide_type: SlideType
    title: str
    subtitle: str = ""
    bullets: List[BulletItem] = field(default_factory=list)
    table: Optional[TableData] = None
    has_image: bool = False  # For future image support


@dataclass
class LayoutInfo:
    """Information about a slide layout"""
    index: int
    name: str
    has_title_placeholder: bool = False
    has_body_placeholder: bool = False
    has_subtitle_placeholder: bool = False
    has_picture_placeholder: bool = False
    placeholder_count: int = 0


# ============================================
# MARKDOWN PARSER
# ============================================

class MarkdownToSlides:
    """Parse markdown and extract slide structure"""
    
    def __init__(self, markdown_text: str):
        self.markdown_text = markdown_text
        self.slides: List[SlideContent] = []
        self.current_h1 = None
        self.current_h2 = None
        self.current_h3 = None
        self.current_slide_content = []
        self._in_html_block = False
        
    def parse(self) -> List[SlideContent]:
        """Parse markdown and return list of slides"""
        lines = self.markdown_text.split('\n')
        i = 0
        
        while i < len(lines):
            line = lines[i]
            
            # Handle multi-line HTML/CSS blocks (e.g., <style>...</style>)
            if self._in_html_block:
                if '</style>' in line or '</script>' in line:
                    self._in_html_block = False
                i += 1
                continue
            
            # Detect start of multi-line HTML blocks
            stripped_line = line.strip()
            if stripped_line.startswith('<style') or stripped_line.startswith('<script'):
                self._in_html_block = True
                # Check if it closes on the same line
                if '</style>' in line or '</script>' in line:
                    self._in_html_block = False
                i += 1
                continue
            
            # Skip standalone HTML comment lines (<!-- ... -->)
            if stripped_line.startswith('<!--') and stripped_line.endswith('-->'):
                i += 1
                continue
            
            # H1 - Title slide
            if line.startswith('# ') and not line.startswith('## '):
                self._flush_current_slide()
                title = line[2:].strip()
                # Check if next line is subtitle (non-empty, non-heading, non-bullet)
                subtitle = ""
                if i + 1 < len(lines):
                    next_line = lines[i + 1].strip()
                    if (next_line 
                            and not next_line.startswith('#')
                            and not next_line.startswith('- ')
                            and not next_line.startswith('* ')
                            and not re.match(r'^\d+\.\s', next_line)
                            and not next_line.startswith('|')):
                        subtitle = next_line
                        i += 1
                
                self.slides.append(SlideContent(
                    slide_type=SlideType.TITLE,
                    title=title,
                    subtitle=subtitle
                ))
                self.current_h1 = title
                i += 1
                continue
            
            # H2 - Section divider
            if line.startswith('## ') and not line.startswith('### '):
                self._flush_current_slide()
                title = line[3:].strip()
                self.slides.append(SlideContent(
                    slide_type=SlideType.SECTION,
                    title=title
                ))
                self.current_h2 = title
                self.current_h3 = None
                i += 1
                continue
            
            # H3 - Content slide
            if line.startswith('### '):
                self._flush_current_slide()
                self.current_h3 = line[4:].strip()
                i += 1
                continue
            
            # Horizontal rule - slide break
            if stripped_line == '---':
                self._flush_current_slide()
                # Keep current H3 for continuation
                i += 1
                continue
            
            # Table detection
            if '|' in line and stripped_line.startswith('|'):
                table_lines = [line]
                i += 1
                # Collect all table lines (must also start with |)
                while i < len(lines) and '|' in lines[i] and lines[i].strip().startswith('|'):
                    table_lines.append(lines[i])
                    i += 1
                
                table = self._parse_table(table_lines)
                if table:
                    self.current_slide_content.append(('table', table))
                continue
            
            # Bullet points (unordered)
            if stripped_line.startswith('- ') or stripped_line.startswith('* '):
                bullet = self._parse_bullet(line)
                if bullet:
                    self.current_slide_content.append(('bullet', bullet))
                i += 1
                continue
            
            # Numbered lists (ordered)
            numbered_match = re.match(r'^(\s*)\d+\.\s+(.*)', line)
            if numbered_match:
                bullet = self._parse_numbered_item(line, numbered_match)
                if bullet:
                    self.current_slide_content.append(('bullet', bullet))
                i += 1
                continue
            
            # Empty lines or other content
            i += 1
        
        # Flush any remaining content
        self._flush_current_slide()
        
        return self.slides
    
    def _flush_current_slide(self):
        """Create a slide from accumulated content"""
        if self.current_h3 and self.current_slide_content:
            bullets = [item[1] for item in self.current_slide_content if item[0] == 'bullet']
            tables = [item[1] for item in self.current_slide_content if item[0] == 'table']
            
            slide = SlideContent(
                slide_type=SlideType.CONTENT,
                title=self.current_h3,
                bullets=bullets,
                table=tables[0] if tables else None
            )
            self.slides.append(slide)
            self.current_slide_content = []
    
    def _parse_bullet(self, line: str) -> Optional[BulletItem]:
        """Parse a bullet point line"""
        # Determine nesting level by leading spaces
        stripped = line.lstrip()
        if not (stripped.startswith('- ') or stripped.startswith('* ')):
            return None
        
        indent = len(line) - len(stripped)
        level = indent // 2  # 2 spaces per level
        
        # Remove bullet marker
        text = stripped[2:].strip()
        
        # Parse inline formatting
        text_runs = self._parse_inline_formatting(text)
        
        return BulletItem(text_runs=text_runs, level=min(level, 2))
    
    def _parse_numbered_item(self, line: str, match: re.Match) -> Optional[BulletItem]:
        """Parse a numbered list item (e.g., '1. Item text')"""
        indent_str = match.group(1)
        text = match.group(2).strip()
        
        indent = len(indent_str)
        level = indent // 2  # 2 spaces per level
        
        # Parse inline formatting
        text_runs = self._parse_inline_formatting(text)
        
        return BulletItem(text_runs=text_runs, level=min(level, 2))
    
    def _parse_inline_formatting(self, text: str) -> List[TextRun]:
        """Parse bold, italic, and links in text"""
        runs = []
        current_pos = 0
        
        # Pattern for **bold**, *italic*, and [link](url)
        pattern = r'(\*\*.*?\*\*|\*.*?\*|\[.*?\]\(.*?\))'
        
        for match in re.finditer(pattern, text):
            # Add text before match
            if match.start() > current_pos:
                runs.append(TextRun(text=text[current_pos:match.start()]))
            
            matched_text = match.group(1)
            
            # Bold
            if matched_text.startswith('**') and matched_text.endswith('**'):
                runs.append(TextRun(text=matched_text[2:-2], bold=True))
            # Italic
            elif matched_text.startswith('*') and matched_text.endswith('*'):
                runs.append(TextRun(text=matched_text[1:-1], italic=True))
            # Link
            elif matched_text.startswith('['):
                link_match = re.match(r'\[(.*?)\]\((.*?)\)', matched_text)
                if link_match:
                    link_text, url = link_match.groups()
                    runs.append(TextRun(text=link_text, link=url))
            
            current_pos = match.end()
        
        # Add remaining text
        if current_pos < len(text):
            runs.append(TextRun(text=text[current_pos:]))
        
        return runs if runs else [TextRun(text=text)]
    
    def _parse_table(self, lines: List[str]) -> Optional[TableData]:
        """Parse markdown table"""
        if len(lines) < 2:
            return None
        
        # Parse header
        header_line = lines[0].strip()
        headers = [cell.strip() for cell in header_line.split('|')[1:-1]]
        
        if not headers:
            return None
        
        num_cols = len(headers)
        
        # Skip separator line (line 1)
        # Parse data rows
        rows = []
        for line in lines[2:]:
            if not line.strip():
                continue
            cells = [cell.strip() for cell in line.split('|')[1:-1]]
            if cells:
                # Normalize row to match header column count
                if len(cells) < num_cols:
                    cells.extend([""] * (num_cols - len(cells)))
                elif len(cells) > num_cols:
                    cells = cells[:num_cols]
                rows.append(cells)
        
        return TableData(headers=headers, rows=rows)


# ============================================
# LAYOUT MANAGER - Intelligent Layout Selection
# ============================================

class LayoutManager:
    """
    Manages intelligent layout selection from PowerPoint templates.
    Analyzes available layouts and maps content types to appropriate layouts.
    """
    
    # Keywords to identify layout types by name
    LAYOUT_KEYWORDS = {
        "title": ["title", "cover", "opening"],
        "section": ["section", "divider", "header"],
        "content": ["content", "body", "text", "bullet"],
        "image": ["image", "picture", "photo", "media"],
        "table": ["table", "comparison"],
        "blank": ["blank", "empty"],
        "two_column": ["two", "column", "side"],
    }
    
    def __init__(self, presentation: 'Presentation'):
        self.prs = presentation
        self.layouts: List[LayoutInfo] = []
        self.layout_map: Dict[str, int] = {}
        self._analyze_layouts()
        self._build_layout_map()
    
    def _analyze_layouts(self):
        """Analyze all layouts in the presentation template"""
        for idx, layout in enumerate(self.prs.slide_layouts):
            info = LayoutInfo(
                index=idx,
                name=layout.name or f"Layout {idx}",
                placeholder_count=0
            )
            
            # Analyze placeholders in the layout
            placeholder_details = []
            for shape in layout.placeholders:
                info.placeholder_count += 1
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                
                # Collect details for debugging
                placeholder_details.append(f"idx={ph_idx}:type={ph_type}")
                
                # Check placeholder types using enum constants
                if ph_type in [PH_TITLE, PH_CENTER_TITLE]:
                    info.has_title_placeholder = True
                elif ph_type in [PH_BODY, PH_OBJECT]:
                    info.has_body_placeholder = True
                elif ph_type == PH_SUBTITLE:
                    info.has_subtitle_placeholder = True
                elif ph_type in [PH_PICTURE, PH_MEDIA_CLIP]:
                    info.has_picture_placeholder = True
                
                # Also check by index for common patterns
                # Index 0 is usually title, Index 1 is usually body/content
                if ph_idx == 1 and not info.has_body_placeholder:
                    # Index 1 often holds content even if type is unusual
                    info.has_body_placeholder = True
            
            self.layouts.append(info)
            logger.debug(
                "  Layout %d: '%s' - Title:%s, Body:%s, Subtitle:%s, Picture:%s [%s]",
                idx, info.name, info.has_title_placeholder,
                info.has_body_placeholder, info.has_subtitle_placeholder,
                info.has_picture_placeholder, ', '.join(placeholder_details)
            )
    
    def _build_layout_map(self):
        """Build a mapping from content types to best layout indices"""
        
        # Find layouts by analyzing names and placeholders
        for info in self.layouts:
            name_lower = info.name.lower()
            
            # Title Slide - exact match for Sand Island template
            if name_lower == "title slide":
                self.layout_map["title_no_image"] = info.index
                self.layout_map["title_with_image"] = info.index  # Same layout for both
            
            # Transition Slide - for section dividers
            if "transition" in name_lower:
                self.layout_map["section"] = info.index
            
            # Content Slide - for regular content
            if name_lower == "content slide":
                self.layout_map["content"] = info.index
            
            # Agenda or Table of Contents - can be used as alternative section or TOC
            if "agenda" in name_lower or "table of contents" in name_lower:
                if "agenda" not in self.layout_map:
                    self.layout_map["agenda"] = info.index
            
            # Generic fallback matching
            # Title/Cover layouts (with subtitle, without image)
            if any(kw in name_lower for kw in ["cover without", "title slide"]):
                if "title_no_image" not in self.layout_map:
                    self.layout_map["title_no_image"] = info.index
            
            # Title/Cover layouts (with image)
            if any(kw in name_lower for kw in ["cover with image", "title with picture"]):
                if info.has_picture_placeholder:
                    if "title_with_image" not in self.layout_map:
                        self.layout_map["title_with_image"] = info.index
            
            # Section header layouts
            if any(kw in name_lower for kw in ["section", "divider"]):
                if "section" not in self.layout_map:
                    self.layout_map["section"] = info.index
            
            # Content layouts (title + content) - relaxed matching
            if "title and content" in name_lower or ("content" in name_lower and "cover" not in name_lower and "table of contents" not in name_lower):
                if info.has_title_placeholder:
                    if "content" not in self.layout_map:
                        self.layout_map["content"] = info.index
            
            # Blank layout
            if "blank" in name_lower or info.placeholder_count == 0:
                if "blank" not in self.layout_map:
                    self.layout_map["blank"] = info.index
        
        # Fallback mappings based on placeholder analysis if names didn't match
        if "title_no_image" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and info.has_subtitle_placeholder and not info.has_picture_placeholder:
                    self.layout_map["title_no_image"] = info.index
                    break
        
        if "content" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and info.has_body_placeholder:
                    self.layout_map["content"] = info.index
                    break
        
        # If still no content layout, use any layout with a title placeholder
        if "content" not in self.layout_map:
            for info in self.layouts:
                if info.has_title_placeholder and not info.has_picture_placeholder:
                    # Avoid cover slides
                    if "cover" not in info.name.lower():
                        self.layout_map["content"] = info.index
                        break
        
        # For section slides, use transition or agenda layout if no dedicated section layout
        if "section" not in self.layout_map:
            if "agenda" in self.layout_map:
                self.layout_map["section"] = self.layout_map["agenda"]
            elif "title_no_image" in self.layout_map:
                self.layout_map["section"] = self.layout_map["title_no_image"]
        
        # Last resort: use indices commonly found in templates
        if "blank" not in self.layout_map:
            try:
                if len(self.prs.slide_layouts) > 6:
                    self.layout_map["blank"] = 6
                else:
                    self.layout_map["blank"] = len(self.prs.slide_layouts) - 1
            except Exception:
                self.layout_map["blank"] = 0
        
        logger.debug("Layout mapping: %s", self.layout_map)
    
    def get_layout_for_slide(self, slide_content: SlideContent) -> Any:
        """
        Intelligently select the best layout for the given slide content.
        
        Args:
            slide_content: The content to be placed on the slide
            
        Returns:
            The appropriate slide layout object
        """
        layout_key = None
        
        if slide_content.slide_type == SlideType.TITLE:
            # Title slides - check for image
            if slide_content.has_image:
                layout_key = "title_with_image"
            else:
                layout_key = "title_no_image"
        
        elif slide_content.slide_type == SlideType.SECTION:
            layout_key = "section"
        
        elif slide_content.slide_type == SlideType.CONTENT:
            # Content slides - could be enhanced for tables, images, etc.
            layout_key = "content"
        
        # Get the layout index, with fallback to blank
        if layout_key and layout_key in self.layout_map:
            idx = self.layout_map[layout_key]
            logger.debug("  -> Using layout '%s' for %s slide", self.layouts[idx].name, slide_content.slide_type.value)
            return self.prs.slide_layouts[idx]
        
        # Fallback to blank layout
        idx = self.layout_map.get("blank", 0)
        logger.debug("  -> Fallback to blank layout for %s slide", slide_content.slide_type.value)
        return self.prs.slide_layouts[idx]
    
    def get_blank_layout(self) -> Any:
        """Get the blank layout for manual slide creation"""
        idx = self.layout_map.get("blank", len(self.prs.slide_layouts) - 1)
        return self.prs.slide_layouts[idx]
    
    def has_usable_layouts(self) -> bool:
        """Check if the template has usable layouts beyond blank"""
        return len(self.layout_map) > 1


# ============================================
# POWERPOINT GENERATOR
# ============================================

def get_color(color_name: str) -> RGBColor:
    """Get RGBColor from color name in config"""
    rgb = CONFIG["colors"][color_name]
    return RGBColor(*rgb)


def _get_fallback_blank_layout(prs: Presentation):
    """Get a blank layout from a presentation without a LayoutManager.
    
    Falls back to index 6 (common blank index) or the last available layout.
    """
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[len(prs.slide_layouts) - 1]


def convert_potx_to_pptx(potx_path: str) -> str:
    """Convert .potx template to .pptx by modifying content types.
    
    Returns:
        Path to a temporary .pptx file. Caller is responsible for cleanup.
    """
    temp_pptx = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
    temp_pptx.close()
    
    # Copy the potx file
    shutil.copy2(potx_path, temp_pptx.name)
    
    # Modify the [Content_Types].xml to change template to presentation
    with zipfile.ZipFile(temp_pptx.name, 'r') as zip_read:
        # Read all files
        file_list = zip_read.namelist()
        
        # Create a new zip with modified content
        temp_zip = tempfile.NamedTemporaryFile(suffix='.zip', delete=False)
        temp_zip.close()
        
        with zipfile.ZipFile(temp_zip.name, 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for file_name in file_list:
                file_data = zip_read.read(file_name)
                
                # Modify [Content_Types].xml
                if file_name == '[Content_Types].xml':
                    file_data = file_data.replace(
                        b'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
                        b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
                    )
                
                zip_write.writestr(file_name, file_data)
    
    # Replace the original temp file
    shutil.move(temp_zip.name, temp_pptx.name)
    
    return temp_pptx.name


def create_clean_template(source_template: str, output_path: str) -> str:
    """
    Create a clean template file (no slides) from an existing template.
    
    This is useful when you have a .pptx with slides and want to create
    a clean template that only contains the layouts/masters.
    
    Args:
        source_template: Path to the source template (can have slides)
        output_path: Path for the clean template output
        
    Returns:
        Path to the created clean template
    """
    logger.info("Creating clean template from: %s", source_template)
    
    # Open the source template
    prs = Presentation(source_template)
    
    # Delete all slides
    num_slides = len(prs.slides)
    if num_slides > 0:
        logger.info("Removing %d slides from template...", num_slides)
        for i in range(num_slides - 1, -1, -1):
            try:
                delete_slide(prs, i)
            except Exception:
                pass
    
    # Save as the clean template
    prs.save(output_path)
    logger.info("Clean template saved to: %s", output_path)
    logger.info("TIP: You can now use this clean template to avoid duplicate warnings.")
    
    return output_path


def delete_slide(prs: Presentation, slide_index: int):
    """
    Delete a slide from the presentation by index.
    Properly removes both the slide reference and the slide part.
    
    Args:
        prs: Presentation object
        slide_index: Index of slide to delete (0-based)
    """
    # Get the slide's relationship ID
    slide_id = prs.slides._sldIdLst[slide_index].rId
    
    # Get the slide part before dropping the relationship
    slide_part = prs.part.related_part(slide_id)
    
    # Remove the slide ID from the presentation's slide list
    del prs.slides._sldIdLst[slide_index]
    
    # Drop the relationship (this disconnects the slide)
    prs.part.drop_rel(slide_id)
    
    # Try to remove the slide part from the package
    try:
        # Access the package and remove the part
        package = prs.part.package
        if hasattr(package, '_parts'):
            # Remove from parts collection if it exists
            parts_to_remove = [p for p in package._parts if p is slide_part]
            for p in parts_to_remove:
                package._parts.remove(p)
    except Exception:
        pass  # If we can't remove the part, it's okay - file will still work


def clear_template_slides(prs: Presentation):
    """
    Remove all existing slides from a template presentation.
    
    This properly deletes slides to avoid duplicate file warnings.
    """
    num_slides = len(prs.slides)
    if num_slides == 0:
        return
    
    logger.info("Clearing %d existing template slides...", num_slides)
    
    # Delete slides from the end to avoid index shifting issues
    deleted_count = 0
    for i in range(num_slides - 1, -1, -1):
        try:
            delete_slide(prs, i)
            deleted_count += 1
        except Exception as e:
            logger.warning("    Could not delete slide %d: %s", i + 1, e)
    
    logger.info("Template slides cleared (%d slides removed).", deleted_count)


def create_presentation(template_path: Optional[str] = None) -> tuple:
    """
    Initialize presentation with template or default settings.
    
    Returns:
        tuple: (Presentation, LayoutManager or None)
    """
    temp_file = None
    layout_manager = None
    
    if template_path and os.path.exists(template_path):
        # Check if it's a .potx file
        if template_path.lower().endswith('.potx'):
            logger.info("Note: .potx template detected. Converting to .pptx format...")
            temp_file = convert_potx_to_pptx(template_path)
            
            try:
                prs = Presentation(temp_file)
                logger.info("Using template: %s", template_path)
            except Exception:
                # Clean up temp file on failure, then re-raise
                if temp_file:
                    try:
                        os.unlink(temp_file)
                    except OSError:
                        pass
                raise
            finally:
                # Clean up temp file after presentation is loaded into memory
                if temp_file:
                    try:
                        os.unlink(temp_file)
                    except OSError:
                        pass
        else:
            prs = Presentation(template_path)
            logger.info("Using template: %s", template_path)
        
        # Clear any existing slides from the template
        if len(prs.slides) > 0:
            clear_template_slides(prs)
        
        # Initialize layout manager for template analysis
        logger.info("Analyzing template layouts...")
        layout_manager = LayoutManager(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(CONFIG["slide_width"])
        prs.slide_height = Inches(CONFIG["slide_height"])
        logger.info("Using default HDOT styling (no template)")
    
    return prs, layout_manager


def _safe_get_first_run(paragraph):
    """Safely get or create the first run of a paragraph after setting text.
    
    After setting `p.text = ...`, python-pptx may or may not create a run object.
    This helper ensures we always have a valid run to format.
    
    Returns:
        The first run of the paragraph.
    """
    if paragraph.runs:
        return paragraph.runs[0]
    # No run was created; add one manually
    run = paragraph.add_run()
    run.text = paragraph.text
    # Clear the direct text so it doesn't duplicate
    paragraph.text = ""
    # Re-fetch since add_run changed the structure
    return paragraph.runs[0] if paragraph.runs else run


def add_title_slide(prs: Presentation, title: str, subtitle: str = "", 
                    layout_manager: Optional[LayoutManager] = None):
    """
    Add a title slide using template layout or manual creation.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text (optional)
        layout_manager: LayoutManager for intelligent layout selection
    """
    # Try to use template layout
    if layout_manager:
        slide_content = SlideContent(slide_type=SlideType.TITLE, title=title, subtitle=subtitle)
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to populate placeholders
        title_set = False
        subtitle_set = False
        
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            
            # TITLE or CENTER_TITLE
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                shape.text = title
                title_set = True
                logger.debug("    -> Title set in placeholder type %s", ph_type)
            
            # SUBTITLE
            elif ph_type == PH_SUBTITLE and subtitle and not subtitle_set:
                shape.text = subtitle
                subtitle_set = True
                logger.debug("    -> Subtitle set in placeholder type %s", ph_type)
        
        if title_set:
            return slide
        
        # If no placeholders found, fall through to manual creation
        logger.debug("    No title placeholders found, using manual layout")
    
    # Fallback: Manual slide creation (original behavior)
    slide_layout = _get_fallback_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, 
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["title_slide_bg"])
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), 
        Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(CONFIG["fonts"]["title"]["size"])
    p.font.bold = CONFIG["fonts"]["title"]["bold"]
    p.font.name = CONFIG["fonts"]["title"]["name"]
    p.font.color.rgb = get_color(CONFIG["color_mapping"]["title_slide_text"])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), 
            Inches(12.333), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.name = CONFIG["fonts"]["title"]["name"]
        p.font.color.rgb = get_color(CONFIG["color_mapping"]["title_slide_text"])
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs: Presentation, title: str,
                      layout_manager: Optional[LayoutManager] = None):
    """
    Add a section divider slide using template layout or manual creation.
    
    Args:
        prs: Presentation object
        title: Section title text
        layout_manager: LayoutManager for intelligent layout selection
    """
    # Try to use template layout
    if layout_manager:
        slide_content = SlideContent(slide_type=SlideType.SECTION, title=title)
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to populate placeholders
        title_set = False
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            # TITLE or CENTER_TITLE
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                shape.text = title
                title_set = True
                logger.debug("    -> Section title set in placeholder type %s", ph_type)
                break
        
        if title_set:
            return slide
        
        # If no placeholders, try adding text to the title area if layout has one
        logger.debug("    No section placeholders found, using manual layout")
    
    # Fallback: Manual slide creation (original behavior)
    slide_layout = _get_fallback_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3), 
        prs.slide_width, Inches(1.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["section_slide_bar"])
    bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), 
        Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(CONFIG["fonts"]["section"]["size"])
    p.font.bold = CONFIG["fonts"]["section"]["bold"]
    p.font.name = CONFIG["fonts"]["section"]["name"]
    p.font.color.rgb = get_color(CONFIG["color_mapping"]["section_slide_text"])
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs: Presentation, slide_content: SlideContent,
                      layout_manager: Optional[LayoutManager] = None):
    """
    Add a content slide using template layout or manual creation.
    
    Args:
        prs: Presentation object
        slide_content: SlideContent with title, bullets, and/or table
        layout_manager: LayoutManager for intelligent layout selection
    """
    # Try to use template layout
    if layout_manager:
        slide_layout = layout_manager.get_layout_for_slide(slide_content)
        slide = prs.slides.add_slide(slide_layout)
        
        # Try to populate placeholders
        title_set = False
        body_set = False
        body_placeholder = None
        
        # First pass: identify placeholders
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            
            # Title placeholder (TITLE or CENTER_TITLE)
            if ph_type in [PH_TITLE, PH_CENTER_TITLE] and not title_set:
                shape.text = slide_content.title
                title_set = True
                logger.debug("    -> Title set in placeholder type %s", ph_type)
            
            # Body/Content placeholder (BODY or OBJECT)
            # OBJECT placeholders can hold text, tables, charts, etc.
            elif ph_type in [PH_BODY, PH_OBJECT] and not body_set:
                body_placeholder = shape
        
        # Second pass: populate body content
        if body_placeholder and slide_content.bullets:
            tf = body_placeholder.text_frame
            
            # Clear any existing placeholder text
            for para in list(tf.paragraphs):
                para.clear()
            
            for i, bullet in enumerate(slide_content.bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.level = bullet.level
                
                # Add text runs with formatting
                for j, run in enumerate(bullet.text_runs):
                    if j == 0:
                        p.text = run.text
                        if p.runs:
                            r = p.runs[0]
                            if run.bold:
                                r.font.bold = True
                            if run.italic:
                                r.font.italic = True
                            if run.link:
                                r.hyperlink.address = run.link
                    else:
                        r = p.add_run()
                        r.text = run.text
                        if run.bold:
                            r.font.bold = True
                        if run.italic:
                            r.font.italic = True
                        if run.link:
                            r.hyperlink.address = run.link
            
            body_set = True
            logger.debug("    -> Body content set in placeholder")
        
        # If we successfully used placeholders for non-table content, return
        if title_set and (body_set or slide_content.table):
            if slide_content.table:
                # Delete the unused body placeholder before adding table
                if body_placeholder and not body_set:
                    try:
                        sp = body_placeholder._element
                        sp.getparent().remove(sp)
                        logger.debug("    -> Removed unused body placeholder")
                    except Exception:
                        pass  # If deletion fails, continue anyway
                
                # Add table separately (tables need precise positioning)
                add_table_to_slide(slide, prs, slide_content.table, use_template_position=True)
            return slide
        
        # If layout had placeholders but content didn't fit, we already have the slide
        # Just need to add content manually
        if title_set:
            if slide_content.table:
                # Delete the unused body placeholder before adding table
                if body_placeholder:
                    try:
                        sp = body_placeholder._element
                        sp.getparent().remove(sp)
                        logger.debug("    -> Removed unused body placeholder")
                    except Exception:
                        pass  # If deletion fails, continue anyway
                
                add_table_to_slide(slide, prs, slide_content.table, use_template_position=True)
            elif slide_content.bullets and not body_set:
                add_bullets_to_slide(slide, slide_content.bullets, use_template_position=True)
            return slide
        
        logger.debug("    No content placeholders found, using manual layout")
    
    # Fallback: Manual slide creation (original behavior)
    slide_layout = _get_fallback_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, 
        prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["content_title_bar"])
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), 
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_content.title
    p.font.size = Pt(CONFIG["fonts"]["slide_title"]["size"])
    p.font.bold = CONFIG["fonts"]["slide_title"]["bold"]
    p.font.name = CONFIG["fonts"]["slide_title"]["name"]
    p.font.color.rgb = get_color(CONFIG["color_mapping"]["content_title_text"])
    
    # Content area
    if slide_content.table:
        add_table_to_slide(slide, prs, slide_content.table)
    elif slide_content.bullets:
        add_bullets_to_slide(slide, slide_content.bullets)
    
    return slide


def add_bullets_to_slide(slide, bullets: List[BulletItem], use_template_position: bool = False):
    """
    Add bullet points to a slide.
    
    Args:
        slide: The slide to add bullets to
        bullets: List of BulletItem objects
        use_template_position: If True, position content lower to avoid template title area
    """
    # Adjust vertical position based on whether we're using a template
    top_position = Inches(1.8) if use_template_position else Inches(1.5)
    
    content_box = slide.shapes.add_textbox(
        Inches(0.5), top_position, 
        Inches(12.333), Inches(5.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Set indentation based on level
        p.level = bullet.level
        
        # Add text runs with formatting
        for j, run in enumerate(bullet.text_runs):
            if j == 0:
                p.text = run.text
                r = _safe_get_first_run(p)
            else:
                r = p.add_run()
                r.text = run.text
            
            r.font.size = Pt(CONFIG["fonts"]["body"]["size"])
            r.font.name = CONFIG["fonts"]["body"]["name"]
            r.font.color.rgb = get_color(CONFIG["color_mapping"]["body_text"])
            
            if run.bold:
                r.font.bold = True
            if run.italic:
                r.font.italic = True
            if run.link:
                r.hyperlink.address = run.link
        
        p.space_after = Pt(8)


def parse_cell_formatting(cell_text: str) -> List[tuple]:
    """
    Parse markdown formatting in table cell text.
    Returns list of (text, is_bold, is_italic) tuples.
    """
    result = []
    text = str(cell_text)
    current_pos = 0
    
    # Pattern for **bold** and *italic*
    pattern = r'(\*\*.*?\*\*|\*.*?\*)'
    
    for match in re.finditer(pattern, text):
        # Add text before match
        if match.start() > current_pos:
            result.append((text[current_pos:match.start()], False, False))
        
        matched_text = match.group(1)
        
        # Bold
        if matched_text.startswith('**') and matched_text.endswith('**'):
            result.append((matched_text[2:-2], True, False))
        # Italic
        elif matched_text.startswith('*') and matched_text.endswith('*'):
            result.append((matched_text[1:-1], False, True))
        
        current_pos = match.end()
    
    # Add remaining text
    if current_pos < len(text):
        result.append((text[current_pos:], False, False))
    
    return result if result else [(text, False, False)]


def set_cell_text_with_formatting(cell, text: str, font_size: int, font_name: str, font_color):
    """
    Set cell text with markdown formatting support (bold/italic).
    """
    parts = parse_cell_formatting(text)
    
    tf = cell.text_frame
    p = tf.paragraphs[0]
    
    for i, (part_text, is_bold, is_italic) in enumerate(parts):
        if i == 0:
            p.text = part_text
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.color.rgb = font_color
                run.font.bold = is_bold
                run.font.italic = is_italic
        else:
            run = p.add_run()
            run.text = part_text
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = font_color
            run.font.bold = is_bold
            run.font.italic = is_italic


def add_table_to_slide(slide, prs: Presentation, table: TableData, use_template_position: bool = False):
    """
    Add a table to a slide with proper sizing and formatting.
    
    Args:
        slide: The slide to add the table to
        prs: Presentation object for dimensions
        table: TableData object with headers and rows
        use_template_position: If True, position content lower to avoid template title area
    """
    num_cols = len(table.headers)
    num_rows = len(table.rows) + 1
    
    # Calculate table dimensions to fit within slide
    left_margin = Inches(0.5)
    right_margin = Inches(0.8)  # Extra space for page number
    
    # Template slides need more bottom margin to clear footer/accent elements
    bottom_margin = Inches(0.9) if use_template_position else Inches(0.4)
    
    # Available width (slide width minus margins)
    slide_width_inches = prs.slide_width.inches
    table_width = Inches(slide_width_inches - left_margin.inches - right_margin.inches)
    
    # Adjust position based on template usage
    top_position = Inches(1.5) if use_template_position else Inches(1.4)
    
    # Available height (slide height minus top position and bottom margin)
    slide_height_inches = prs.slide_height.inches
    available_height = slide_height_inches - top_position.inches - bottom_margin.inches
    
    # Calculate row height based on number of rows
    # Use tighter minimum for large tables to ensure they fit
    min_row_height = 0.28 if num_rows > 8 else 0.35  # inches
    calculated_row_height = available_height / num_rows
    row_height = max(min_row_height, calculated_row_height)
    
    # Clamp table height to available space
    table_height = Inches(min(available_height, row_height * num_rows))
    
    pptx_table = slide.shapes.add_table(
        num_rows, num_cols, 
        left_margin, top_position, 
        table_width, table_height
    ).table
    
    # Set column widths evenly
    col_width = int(table_width / num_cols)
    for i in range(num_cols):
        pptx_table.columns[i].width = col_width
    
    # For large tables, reduce font sizes and cell margins to fit
    is_large_table = num_rows > 8
    header_font_size = CONFIG["fonts"]["table_header"]["size"] - (2 if is_large_table else 0)
    body_font_size = CONFIG["fonts"]["table_body"]["size"] - (1 if is_large_table else 0)
    
    # Cell margin settings (reduce for large tables)
    cell_margin_top = Inches(0.02) if is_large_table else Inches(0.05)
    cell_margin_bottom = Inches(0.02) if is_large_table else Inches(0.05)
    cell_margin_left = Inches(0.05)
    cell_margin_right = Inches(0.05)
    
    # Header row
    for i, header in enumerate(table.headers):
        cell = pptx_table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["table_header_bg"])
        
        # Reduce cell margins for tighter fit
        cell.margin_top = cell_margin_top
        cell.margin_bottom = cell_margin_bottom
        cell.margin_left = cell_margin_left
        cell.margin_right = cell_margin_right
        
        p = cell.text_frame.paragraphs[0]
        p.font.bold = CONFIG["fonts"]["table_header"]["bold"]
        p.font.size = Pt(header_font_size)
        p.font.name = CONFIG["fonts"]["table_header"]["name"]
        p.font.color.rgb = get_color(CONFIG["color_mapping"]["table_header_text"])
        p.alignment = PP_ALIGN.CENTER
        # Vertical alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows with formatting support
    for row_idx, row in enumerate(table.rows):
        for col_idx in range(num_cols):
            cell = pptx_table.cell(row_idx + 1, col_idx)
            cell_text = row[col_idx] if col_idx < len(row) else ""
            
            # Reduce cell margins for tighter fit
            cell.margin_top = cell_margin_top
            cell.margin_bottom = cell_margin_bottom
            cell.margin_left = cell_margin_left
            cell.margin_right = cell_margin_right
            
            # Use formatted text setting to handle **bold** markers
            set_cell_text_with_formatting(
                cell, 
                cell_text,
                body_font_size,
                CONFIG["fonts"]["table_body"]["name"],
                get_color(CONFIG["color_mapping"]["body_text"])
            )
            
            # Vertical alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = get_color(CONFIG["color_mapping"]["table_alt_row"])


# ============================================
# MAIN CONVERSION FUNCTION
# ============================================

def convert_markdown_to_pptx(
    markdown_path: str, 
    output_path: Optional[str] = None,
    template_path: Optional[str] = None
) -> str:
    """
    Convert markdown file to PowerPoint presentation
    
    Args:
        markdown_path: Path to input markdown file
        output_path: Path to output PPTX file (optional)
        template_path: Path to PowerPoint template file (optional)
    
    Returns:
        Path to generated PPTX file
    """
    # Read markdown file
    with open(markdown_path, 'r', encoding='utf-8') as f:
        markdown_text = f.read()
    
    # Parse markdown
    logger.info("Parsing markdown...")
    parser = MarkdownToSlides(markdown_text)
    slides = parser.parse()
    logger.info("Found %d slides", len(slides))
    
    # Create presentation with layout manager
    logger.info("Creating presentation...")
    prs, layout_manager = create_presentation(template_path)
    
    if layout_manager:
        logger.info("Intelligent layout selection enabled with %d layouts available", len(layout_manager.layouts))
    else:
        logger.info("Using manual slide creation (no template)")
    
    # Generate slides with intelligent layout selection
    logger.info("Generating slides...")
    for i, slide_content in enumerate(slides):
        title_display = (slide_content.title[:50] + "...") if len(slide_content.title) > 50 else slide_content.title
        logger.info("Slide %d: %s - '%s'", i + 1, slide_content.slide_type.value, title_display)
        
        if slide_content.slide_type == SlideType.TITLE:
            add_title_slide(prs, slide_content.title, slide_content.subtitle, layout_manager)
        elif slide_content.slide_type == SlideType.SECTION:
            add_section_slide(prs, slide_content.title, layout_manager)
        elif slide_content.slide_type == SlideType.CONTENT:
            add_content_slide(prs, slide_content, layout_manager)
    
    # Determine output path
    if not output_path:
        base = os.path.splitext(markdown_path)[0]
        output_path = f"{base}.pptx"
    
    # Save presentation
    prs.save(output_path)
    logger.info("=" * 50)
    logger.info("Presentation saved to: %s", output_path)
    logger.info("Total slides: %d", len(prs.slides))
    
    return output_path


# ============================================
# COMMAND LINE INTERFACE
# ============================================

def main():
    parser = argparse.ArgumentParser(
        description="Convert markdown files to PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python md_to_pptx_converter.py document.md
  python md_to_pptx_converter.py document.md --template template.pptx
  python md_to_pptx_converter.py document.md -o output.pptx
  
  # Create a clean template (no slides) from an existing one:
  python md_to_pptx_converter.py --clean-template source.pptx -o clean_template.pptx
        """
    )
    
    parser.add_argument(
        'markdown_file',
        nargs='?',  # Make optional for --clean-template mode
        help='Path to input markdown file'
    )
    
    parser.add_argument(
        '-o', '--output',
        help='Path to output PPTX file (default: same name as input with .pptx extension)'
    )
    
    parser.add_argument(
        '-t', '--template',
        help='Path to PowerPoint template file (.pptx or .potx)'
    )
    
    parser.add_argument(
        '--clean-template',
        metavar='SOURCE',
        help='Create a clean template (no slides) from SOURCE template file'
    )
    
    parser.add_argument(
        '-v', '--verbose',
        action='store_true',
        help='Enable verbose/debug output'
    )
    
    args = parser.parse_args()
    
    # Configure logging based on verbosity
    log_level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(
        level=log_level,
        format='%(message)s'
    )
    
    # Handle --clean-template mode
    if args.clean_template:
        if not os.path.exists(args.clean_template):
            logger.error("Error: Source template not found: %s", args.clean_template)
            return 1
        
        output_path = args.output or args.clean_template.replace('.pptx', '_clean.pptx')
        try:
            create_clean_template(args.clean_template, output_path)
            return 0
        except Exception as e:
            logger.error("Error creating clean template: %s", e)
            return 1
    
    # Normal conversion mode requires markdown file
    if not args.markdown_file:
        parser.print_help()
        return 1
    
    # Validate input file
    if not os.path.exists(args.markdown_file):
        logger.error("Error: Input file not found: %s", args.markdown_file)
        return 1
    
    # Validate template if provided
    if args.template and not os.path.exists(args.template):
        logger.warning("Warning: Template file not found: %s", args.template)
        logger.warning("Proceeding with default styling...")
        args.template = None
    
    # Convert
    try:
        output_path = convert_markdown_to_pptx(
            args.markdown_file,
            args.output,
            args.template
        )
        logger.info("Success! Created: %s", output_path)
        return 0
    except Exception as e:
        logger.error("Error: %s", e)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())
